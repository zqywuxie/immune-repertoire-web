"""
PPT热图对比服务
支持扫描多个方法的热图文件夹，并生成对比布局的PPT
利用PPT模板中的热图位置信息，在原有位置上进行多项目对比
"""

import os
import re
import base64
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import logging

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import io

logger = logging.getLogger(__name__)


class PPTComparisonService:
    """PPT热图对比服务"""
    
    # 支持的图片格式
    SUPPORTED_IMAGE_FORMATS = ['.png', '.jpg', '.jpeg']
    
    # 指标名称映射（与ppt_heatmap_service保持一致）
    METRIC_PATTERNS = {
        'expression': 'Expression Sharing',
        'morisita': 'Morisita-Horn Index',
        'cdr3': 'Unique CDR3 Sharing',
        'ucdr3': 'Unique CDR3 Sharing',
        'r2_inner': 'R² Inner',
        'r2_outer': 'R² Outer',
        'sorensen': 'Sorensen-Dice Index'
    }
    
    # 指标到幻灯片类型的映射
    # Slide 1: Expression, R² Outer, R² Inner
    # Slide 2: Morisita-Horn, uCDR3, Sorensen
    METRIC_TO_SLIDE_TYPE = {
        'expression': 'expression_r2',
        'r2_outer': 'expression_r2',
        'r2_inner': 'expression_r2',
        'morisita': 'morisita_sorensen',
        'morisita_horn': 'morisita_sorensen',
        'ucdr3': 'morisita_sorensen',
        'cdr3': 'morisita_sorensen',
        'sorensen': 'morisita_sorensen'
    }
    
    # 每种幻灯片类型的指标顺序（从左到右）
    SLIDE_TYPE_METRICS = {
        'expression_r2': ['expression', 'r2_outer', 'r2_inner'],
        'morisita_sorensen': ['morisita', 'ucdr3', 'sorensen']
    }
    
    def __init__(self):
        """初始化服务"""
        self.ppt_service = None  # 延迟导入避免循环依赖
    
    def scan_heatmap_folder(self, folder_path: str) -> List[Dict]:
        """
        扫描文件夹中的热图文件（支持子目录结构）
        
        Args:
            folder_path: 文件夹路径
            
        Returns:
            热图列表，每个热图包含：
            - filename: 文件名
            - filepath: 完整路径
            - metric: 指标类型
            - metric_display: 指标显示名称
            - image_data: Base64编码的图片数据
        """
        folder = Path(folder_path)
        
        if not folder.exists() or not folder.is_dir():
            raise ValueError(f"路径不存在或不是目录: {folder_path}")
        
        heatmaps = []
        
        # 支持的链类型子目录
        chain_types = ['IGH', 'IGK', 'IGL', 'TRA', 'TRB', 'TRD', 'TRG']
        
        # 检查是否有链类型子目录
        has_chain_subdirs = any((folder / chain).exists() and (folder / chain).is_dir() 
                                for chain in chain_types)
        
        if has_chain_subdirs:
            # 扫描每个链类型子目录
            for chain in chain_types:
                chain_dir = folder / chain
                if not chain_dir.exists() or not chain_dir.is_dir():
                    continue
                
                # 检查是否有 heatmaps 子目录（某些格式会有）
                search_dirs = []
                heatmaps_subdir = chain_dir / 'heatmaps'
                if heatmaps_subdir.exists() and heatmaps_subdir.is_dir():
                    search_dirs.append(heatmaps_subdir)
                search_dirs.append(chain_dir)
                
                # 扫描每个搜索目录
                for search_dir in search_dirs:
                    heatmaps.extend(self._scan_directory_for_heatmaps(search_dir))
        else:
            # 直接扫描根目录
            heatmaps = self._scan_directory_for_heatmaps(folder)
        
        # 按指标排序
        heatmaps.sort(key=lambda x: list(self.METRIC_PATTERNS.keys()).index(x['metric']) 
                      if x['metric'] in self.METRIC_PATTERNS else 999)
        
        logger.info(f"在 {folder_path} 中找到 {len(heatmaps)} 个热图")
        return heatmaps
    
    def _scan_directory_for_heatmaps(self, directory: Path) -> List[Dict]:
        """
        扫描单个目录中的热图文件
        
        Args:
            directory: 要扫描的目录
            
        Returns:
            热图列表
        """
        heatmaps = []
        
        # 遍历目录中的所有图片文件
        for file_path in directory.iterdir():
            if not file_path.is_file():
                continue
            
            # 检查文件扩展名
            if file_path.suffix.lower() not in self.SUPPORTED_IMAGE_FORMATS:
                continue
            
            # 识别指标类型
            metric = self._identify_metric(file_path.name)
            if not metric:
                logger.debug(f"无法识别指标类型: {file_path.name}")
                continue
            
            # 读取图片并转换为Base64
            try:
                with open(file_path, 'rb') as f:
                    image_data = base64.b64encode(f.read()).decode('utf-8')
                
                heatmaps.append({
                    'filename': file_path.name,
                    'filepath': str(file_path),
                    'metric': metric,
                    'metric_display': self.METRIC_PATTERNS.get(metric, metric),
                    'image_data': image_data
                })
                logger.debug(f"找到热图: {file_path.name} -> {metric}")
            except Exception as e:
                logger.error(f"读取图片失败 {file_path}: {e}")
                continue
        
        return heatmaps
    
    def _identify_metric(self, filename: str) -> Optional[str]:
        """
        从文件名识别指标类型
        
        Args:
            filename: 文件名
            
        Returns:
            指标类型，如 'expression', 'morisita' 等
        """
        filename_lower = filename.lower()
        
        # 移除常见后缀以提高匹配准确性
        filename_lower = filename_lower.replace('_heatmap', '').replace('_sharing', '').replace('.png', '').replace('.jpg', '').replace('.jpeg', '')
        
        # 检查每个指标模式（按优先级顺序）
        # R² Inner - 需要先检查，避免被 r2_outer 误匹配
        if 'r2_inner' in filename_lower or 'r²_inner' in filename_lower or 'r2inner' in filename_lower:
            return 'r2_inner'
        # R² Outer
        elif 'r2_outer' in filename_lower or 'r²_outer' in filename_lower or 'r2outer' in filename_lower:
            return 'r2_outer'
        # Expression
        elif 'expression' in filename_lower or 'expr' in filename_lower:
            return 'expression'
        # Morisita-Horn
        elif 'morisita' in filename_lower or 'horn' in filename_lower:
            return 'morisita'
        # CDR3 Sharing (支持 ucdr3 和 cdr3)
        elif 'ucdr3' in filename_lower or ('cdr3' in filename_lower):
            return 'cdr3'
        # Sorensen-Dice
        elif 'sorensen' in filename_lower or 'dice' in filename_lower:
            return 'sorensen'
        
        return None
    
    def generate_comparison_ppt(
        self,
        template_ppt_path: str,
        methods: List[Dict],
        output_path: str,
        layout_mode: str = 'auto'
    ) -> str:
        """
        生成对比布局的PPT（利用模板中的热图位置信息）
        
        Args:
            template_ppt_path: PPT模板路径
            methods: 方法列表，每个方法包含name和heatmaps
            output_path: 输出路径
            layout_mode: 布局模式
                - 'auto': 自动选择（2-3个项目用单行，4+个项目用网格）
                - 'single_row': 强制使用单行布局
                - 'grid': 强制使用网格布局
            
        Returns:
            生成的PPT文件路径
        """
        logger.info(f"开始生成对比PPT，方法数量: {len(methods)}，布局模式: {layout_mode}")
        
        # 延迟导入避免循环依赖
        from services.ppt_heatmap_service import PPTImageService
        
        # 使用PPT热图服务分析模板
        ppt_service = PPTImageService()
        ppt_service.load_presentation(template_ppt_path)
        slide_info_list = ppt_service.analyze_slides()
        
        logger.info(f"分析PPT模板，找到 {len(slide_info_list)} 个热图页")
        
        # 按链类型和幻灯片类型组织热图页
        slides_by_chain_and_type = {}
        for slide_info in slide_info_list:
            if slide_info.image_type == 'sharing_analysis':
                chain = slide_info.chain_type
                slide_type = slide_info.metric_type
                key = (chain, slide_type)
                slides_by_chain_and_type[key] = slide_info
                logger.info(f"找到热图页: {chain} - {slide_type} (幻灯片 {slide_info.slide_index})")
        
        # 组织方法的热图数据
        methods_heatmaps = self._organize_methods_heatmaps(methods)
        
        # 确定使用的布局方法
        n_methods = len(methods)
        if layout_mode == 'auto':
            use_grid = n_methods > 3
        elif layout_mode == 'grid':
            use_grid = True
        else:  # 'single_row'
            use_grid = False
        
        logger.info(f"使用布局方法: {'网格布局' if use_grid else '单行布局'}")
        
        # 处理每个热图页
        processed_count = 0
        for (chain, slide_type), slide_info in slides_by_chain_and_type.items():
            slide = ppt_service.presentation.slides[slide_info.slide_index]
            
            # 获取该页的指标列表
            metrics = self.SLIDE_TYPE_METRICS.get(slide_type, [])
            
            # 清空原有图片
            self._clear_slide_images(slide)
            
            # 在原有位置上添加多项目对比
            if use_grid:
                self._add_comparison_in_positions_grid(
                    slide, slide_info, chain, metrics, methods_heatmaps
                )
            else:
                self._add_comparison_in_positions(
                    slide, slide_info, chain, metrics, methods_heatmaps
                )
            
            processed_count += 1
            logger.info(f"处理完成: {chain} - {slide_type} ({processed_count}/{len(slides_by_chain_and_type)})")
        
        # 保存PPT
        ppt_service.save_presentation(output_path)
        logger.info(f"对比PPT已保存: {output_path}，共处理 {processed_count} 个热图页")
        
        return output_path
    
    def _organize_methods_heatmaps(self, methods: List[Dict]) -> Dict:
        """
        组织方法的热图数据，按链类型和指标分类
        
        Returns:
            {
                'IGH': {
                    'expression': [
                        {'method_name': '项目A', 'image_data': '...'},
                        {'method_name': '项目B', 'image_data': '...'}
                    ],
                    'r2_outer': [...]
                }
            }
        """
        # 指标别名映射（统一到标准名称）
        metric_aliases = {
            'cdr3': 'ucdr3',
            'unique_cdr3': 'ucdr3',
            'morisita_horn': 'morisita',
            'sorensen_dice': 'sorensen'
        }
        
        organized = {}
        
        for method in methods:
            method_name = method['name']
            
            for heatmap in method['heatmaps']:
                metric = heatmap['metric'].lower()
                
                # 统一指标名称
                metric = metric_aliases.get(metric, metric)
                
                # 从文件名或路径推断链类型（如果有filepath）
                chain = 'IGH'  # 默认
                if 'filepath' in heatmap:
                    filepath = heatmap['filepath']
                    # 尝试从路径中提取链类型
                    for chain_type in ['IGH', 'IGK', 'IGL', 'TRA', 'TRB', 'TRD', 'TRG']:
                        if chain_type in filepath.upper():
                            chain = chain_type
                            break
                
                # 初始化结构
                if chain not in organized:
                    organized[chain] = {}
                if metric not in organized[chain]:
                    organized[chain][metric] = []
                
                # 添加热图数据
                organized[chain][metric].append({
                    'method_name': method_name,
                    'image_data': heatmap.get('image_data', ''),
                    'filepath': heatmap.get('filepath', '')
                })
                
                logger.info(f"组织热图数据: {chain} - {metric} - {method_name}")
        
        return organized
    
    def _add_comparison_in_positions(
        self, 
        slide, 
        slide_info, 
        chain: str, 
        metrics: List[str], 
        methods_heatmaps: Dict
    ):
        """
        在原有热图位置上添加多项目对比（单行布局）
        
        Args:
            slide: 幻灯片对象
            slide_info: 幻灯片信息（包含原有图片位置）
            chain: 链类型
            metrics: 指标列表（按顺序）
            methods_heatmaps: 组织好的方法热图数据
        """
        image_positions = slide_info.image_positions
        n_positions = len(image_positions)
        
        logger.info(f"在 {chain} 页面添加对比，共 {n_positions} 个位置，指标: {metrics}")
        
        # 获取该链的热图数据
        chain_heatmaps = methods_heatmaps.get(chain, {})
        
        # 处理每个位置（对应一个指标）
        for pos_idx, metric in enumerate(metrics[:n_positions]):
            if pos_idx >= len(image_positions):
                break
            
            # 获取原始图片位置
            orig_pos = image_positions[pos_idx]
            orig_left = orig_pos['left_inches']
            orig_top = orig_pos['top_inches']
            orig_width = orig_pos['width_inches']
            orig_height = orig_pos['height_inches']
            
            # 获取该指标的所有项目热图
            metric_heatmaps = chain_heatmaps.get(metric, [])
            
            if not metric_heatmaps:
                logger.warning(f"指标 {metric} 没有热图数据")
                continue
            
            n_methods = len(metric_heatmaps)
            
            # 重新设计布局：每个项目的图片+项目名作为一个整体单元
            label_height = 0.25  # 项目名称高度（减小）
            unit_gap = 0.05  # 单元之间的间隙（减小）
            
            # 计算每个单元的可用高度（图片+标签）
            total_gap = unit_gap * (n_methods - 1)
            available_height = orig_height - total_gap
            unit_height = available_height / n_methods  # 每个单元的总高度
            
            # 图片高度 = 单元高度 - 标签高度
            img_height_per_unit = unit_height - label_height
            
            logger.info(f"布局参数: {n_methods}个项目, 可用高度={available_height:.2f}, "
                       f"单元高度={unit_height:.2f}, 图片高度={img_height_per_unit:.2f}, 标签高度={label_height:.2f}")
            
            # 当前垂直位置（从热图位置顶部开始）
            current_top = orig_top
            
            # 添加每个项目的热图
            for method_idx, heatmap_data in enumerate(metric_heatmaps):
                method_name = heatmap_data['method_name']
                image_data = heatmap_data['image_data']
                
                if not image_data:
                    logger.warning(f"项目 {method_name} 的 {metric} 热图数据为空")
                    continue
                
                try:
                    # 解码图片
                    image_bytes = base64.b64decode(image_data)
                    image_stream = io.BytesIO(image_bytes)
                    
                    # 获取原始尺寸
                    img = Image.open(io.BytesIO(image_bytes))
                    orig_img_width, orig_img_height = img.size
                    aspect_ratio = orig_img_height / orig_img_width
                    
                    # 计算实际图片尺寸（保持纵横比）
                    actual_img_width = orig_width * 0.96  # 使用96%的宽度，留出边距
                    actual_img_height = actual_img_width * aspect_ratio
                    
                    # 如果高度超限，先按高度缩放
                    if actual_img_height > img_height_per_unit:
                        actual_img_height = img_height_per_unit
                        actual_img_width = actual_img_height / aspect_ratio
                    
                    # 在限制范围内增大50%
                    if actual_img_height * 1.5 <= img_height_per_unit:
                        actual_img_height = actual_img_height * 1.5
                        actual_img_width = actual_img_width * 1.5
                    else:
                        # 如果放大50%会超限，则使用最大允许尺寸
                        actual_img_height = img_height_per_unit
                        actual_img_width = actual_img_height / aspect_ratio
                    
                    # 图片位置：从当前垂直位置开始，水平居中
                    img_left = orig_left + (orig_width - actual_img_width) / 2
                    img_top = current_top
                    
                    # 重置stream
                    image_stream.seek(0)
                    
                    # 添加图片
                    pic = slide.shapes.add_picture(
                        image_stream,
                        Inches(img_left), Inches(img_top),
                        width=Inches(actual_img_width),
                        height=Inches(actual_img_height)
                    )
                    
                    # 添加边框
                    line = pic.line
                    line.color.rgb = RGBColor(200, 200, 200)
                    line.width = Pt(0.75)
                    
                    # 项目名称标签：紧跟在图片下方
                    label_top = img_top + actual_img_height + 0.05
                    
                    # 添加项目名称标签（使用完整宽度）
                    label_box = slide.shapes.add_textbox(
                        Inches(orig_left), Inches(label_top),
                        Inches(orig_width), Inches(label_height)
                    )
                    label_frame = label_box.text_frame
                    label_frame.text = method_name
                    label_frame.paragraphs[0].font.size = Pt(11)
                    label_frame.paragraphs[0].font.bold = True
                    label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # 设置文本框背景为透明
                    label_box.fill.background()
                    
                    # 更新当前垂直位置：移动到下一个单元的起始位置
                    # 当前单元结束位置 = 图片顶部 + 图片高度 + 标签高度 + 单元间隙
                    current_top = img_top + actual_img_height + label_height + unit_gap
                    
                    logger.info(f"✓ 添加单元{method_idx+1}: {chain} - {metric} - {method_name} "
                               f"(图片: top={img_top:.2f}, 尺寸={actual_img_width:.2f}x{actual_img_height:.2f}, "
                               f"标签: top={label_top:.2f}, 下一单元起始={current_top:.2f})")
                    
                except Exception as e:
                    logger.error(f"添加图片失败 {chain} - {metric} - {method_name}: {e}")
                    import traceback
                    logger.error(traceback.format_exc())
                    continue
    
    def _add_comparison_in_positions_grid(
        self, 
        slide, 
        slide_info, 
        chain: str, 
        metrics: List[str], 
        methods_heatmaps: Dict
    ):
        """
        在原有热图位置上添加多项目对比（网格布局）
        
        Args:
            slide: 幻灯片对象
            slide_info: 幻灯片信息（包含原有图片位置）
            chain: 链类型
            metrics: 指标列表（按顺序）
            methods_heatmaps: 组织好的方法热图数据
        """
        image_positions = slide_info.image_positions
        n_positions = len(image_positions)
        
        logger.info(f"在 {chain} 页面添加网格对比，共 {n_positions} 个位置，指标: {metrics}")
        
        # 获取该链的热图数据
        chain_heatmaps = methods_heatmaps.get(chain, {})
        
        # 处理每个位置（对应一个指标）
        for pos_idx, metric in enumerate(metrics[:n_positions]):
            if pos_idx >= len(image_positions):
                break
            
            # 获取原始图片位置
            orig_pos = image_positions[pos_idx]
            orig_left = orig_pos['left_inches']
            orig_top = orig_pos['top_inches']
            orig_width = orig_pos['width_inches']
            orig_height = orig_pos['height_inches']
            
            # 获取该指标的所有项目热图
            metric_heatmaps = chain_heatmaps.get(metric, [])
            
            if not metric_heatmaps:
                logger.warning(f"指标 {metric} 没有热图数据")
                continue
            
            n_methods = len(metric_heatmaps)
            
            # 计算网格布局
            if n_methods <= 2:
                n_rows, n_cols = 1, n_methods
            elif n_methods <= 4:
                n_rows, n_cols = 2, 2
            elif n_methods <= 6:
                n_rows, n_cols = 2, 3
            else:
                n_rows, n_cols = 3, 3
            
            # 网格布局：每个单元格内图片+项目名作为整体
            h_gap = 0.08  # 水平间隙（减小）
            v_gap = 0.08  # 垂直间隙（减小）
            label_height = 0.25  # 项目名称标签高度（减小）
            
            # 计算单元格尺寸
            cell_width = (orig_width - h_gap * (n_cols - 1)) / n_cols
            cell_height = (orig_height - v_gap * (n_rows - 1)) / n_rows
            
            # 图片可用高度 = 单元格高度 - 标签高度 - 间隙
            img_height_per_cell = cell_height - label_height - 0.05  # 减小预留间隙
            
            logger.info(f"网格布局参数: {n_methods}个项目, {n_rows}行x{n_cols}列, "
                       f"单元格={cell_width:.2f}x{cell_height:.2f}, 图片高度={img_height_per_cell:.2f}")
            
            # 添加每个项目的热图
            for method_idx, heatmap_data in enumerate(metric_heatmaps):
                if method_idx >= n_rows * n_cols:
                    break
                
                method_name = heatmap_data['method_name']
                image_data = heatmap_data['image_data']
                
                if not image_data:
                    logger.warning(f"项目 {method_name} 的 {metric} 热图数据为空")
                    continue
                
                try:
                    # 计算网格位置
                    row = method_idx // n_cols
                    col = method_idx % n_cols
                    
                    cell_left = orig_left + col * (cell_width + h_gap)
                    cell_top = orig_top + row * (cell_height + v_gap)
                    
                    # 解码图片
                    image_bytes = base64.b64decode(image_data)
                    image_stream = io.BytesIO(image_bytes)
                    
                    # 获取原始尺寸
                    img = Image.open(io.BytesIO(image_bytes))
                    orig_img_width, orig_img_height = img.size
                    aspect_ratio = orig_img_height / orig_img_width
                    
                    # 计算实际图片尺寸（保持纵横比）
                    actual_img_width = cell_width * 0.92  # 使用92%的宽度
                    actual_img_height = actual_img_width * aspect_ratio
                    
                    # 如果高度超限，先按高度缩放
                    if actual_img_height > img_height_per_cell:
                        actual_img_height = img_height_per_cell
                        actual_img_width = actual_img_height / aspect_ratio
                    
                    # 在限制范围内增大50%
                    if actual_img_height * 1.5 <= img_height_per_cell:
                        actual_img_height = actual_img_height * 1.5
                        actual_img_width = actual_img_width * 1.5
                    else:
                        # 如果放大50%会超限，则使用最大允许尺寸
                        actual_img_height = img_height_per_cell
                        actual_img_width = actual_img_height / aspect_ratio
                    
                    # 图片位置：在单元格内水平居中
                    img_left = cell_left + (cell_width - actual_img_width) / 2
                    img_top = cell_top
                    
                    # 重置stream
                    image_stream.seek(0)
                    
                    # 添加图片
                    pic = slide.shapes.add_picture(
                        image_stream,
                        Inches(img_left), Inches(img_top),
                        width=Inches(actual_img_width),
                        height=Inches(actual_img_height)
                    )
                    
                    # 添加边框
                    line = pic.line
                    line.color.rgb = RGBColor(180, 180, 180)
                    line.width = Pt(0.5)
                    
                    # 项目名称标签：紧跟在图片下方
                    label_top = img_top + actual_img_height + 0.05
                    
                    # 添加项目名称标签（使用单元格宽度）
                    label_box = slide.shapes.add_textbox(
                        Inches(cell_left), Inches(label_top),
                        Inches(cell_width), Inches(label_height)
                    )
                    label_frame = label_box.text_frame
                    label_frame.text = method_name
                    label_frame.paragraphs[0].font.size = Pt(10)
                    label_frame.paragraphs[0].font.bold = True
                    label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # 设置文本框背景为透明
                    label_box.fill.background()
                    
                    logger.info(f"✓ 添加单元(网格): {chain} - {metric} - {method_name} "
                               f"(第{row+1}行第{col+1}列, 图片: top={img_top:.2f}, "
                               f"尺寸={actual_img_width:.2f}x{actual_img_height:.2f}, "
                               f"标签: top={label_top:.2f})")
                    
                except Exception as e:
                    logger.error(f"添加图片失败 {chain} - {metric} - {method_name}: {e}")
                    import traceback
                    logger.error(traceback.format_exc())
                    continue
    
    def _find_or_create_slide(self, prs: Presentation, metric: str):
        """
        查找或创建包含指定指标的幻灯片
        
        Args:
            prs: Presentation对象
            metric: 指标类型
            
        Returns:
            幻灯片对象
        """
        metric_display = self.METRIC_PATTERNS.get(metric, metric)
        
        # 查找包含该指标关键词的幻灯片
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    if metric in shape.text.lower() or metric_display.lower() in shape.text.lower():
                        logger.info(f"找到现有幻灯片: {metric_display}")
                        return slide
        
        # 如果没找到，创建新幻灯片
        logger.info(f"创建新幻灯片: {metric_display}")
        blank_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(blank_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = metric_display
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _clear_slide_images(self, slide):
        """清空幻灯片中的图片"""
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE类型
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
    
    def _add_comparison_layout(self, slide, metric: str, methods: List[Dict]):
        """
        在幻灯片中添加对比布局
        
        优化要点：
        1. 项目名称标签放在图片下方
        2. 保持图片原始纵横比
        3. 自动调整布局适应不同数量的项目
        4. 添加图片边框以区分不同项目
        
        Args:
            slide: 幻灯片对象
            metric: 指标类型
            methods: 方法列表
        """
        n_methods = len(methods)
        
        # 计算布局参数
        # 幻灯片尺寸：宽10英寸，高7.5英寸
        slide_width = 10.0
        slide_height = 7.5
        
        # 留出空间：顶部1.2英寸（标题+间距），底部0.8英寸（标签+间距），左右各0.4英寸
        margin_top = 1.2
        margin_bottom = 0.8
        margin_left = 0.4
        margin_right = 0.4
        
        # 可用空间
        available_width = slide_width - margin_left - margin_right
        available_height = slide_height - margin_top - margin_bottom
        
        # 标签高度
        label_height = 0.4
        
        # 图片间隙
        gap = 0.3 if n_methods <= 3 else 0.2
        
        # 计算每个图片的宽度
        img_width = (available_width - gap * (n_methods - 1)) / n_methods
        
        # 图片最大高度（留出标签空间）
        max_img_height = available_height - label_height - 0.1  # 0.1英寸间隙
        
        # 添加每个方法的热图
        for i, method in enumerate(methods):
            # 查找该方法的对应指标热图
            heatmap = next((hm for hm in method['heatmaps'] if hm['metric'] == metric), None)
            
            # 计算当前图片的左边距
            left = Inches(margin_left + i * (img_width + gap))
            
            if not heatmap:
                # 如果该方法没有此指标，显示占位符
                top = Inches(margin_top)
                
                placeholder = slide.shapes.add_textbox(
                    left, top,
                    Inches(img_width), Inches(max_img_height)
                )
                text_frame = placeholder.text_frame
                text_frame.text = f"{method['name']}\n\n无此指标数据"
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame.paragraphs[0].font.size = Pt(16)
                text_frame.paragraphs[0].font.color.rgb = (128, 128, 128)
                continue
            
            # 添加图片
            try:
                # 将Base64解码为图片
                image_bytes = base64.b64decode(heatmap['image_data'])
                image_stream = io.BytesIO(image_bytes)
                
                # 获取图片原始尺寸以保持纵横比
                img = Image.open(io.BytesIO(image_bytes))
                orig_width, orig_height = img.size
                aspect_ratio = orig_height / orig_width
                
                # 计算实际图片尺寸（保持纵横比）
                actual_img_width = img_width
                actual_img_height = img_width * aspect_ratio
                
                # 如果高度超过最大高度，按高度缩放
                if actual_img_height > max_img_height:
                    actual_img_height = max_img_height
                    actual_img_width = max_img_height / aspect_ratio
                
                # 居中对齐（如果图片宽度小于分配宽度）
                img_left = left
                if actual_img_width < img_width:
                    img_left = Inches(margin_left + i * (img_width + gap) + (img_width - actual_img_width) / 2)
                
                # 图片顶部位置
                top = Inches(margin_top)
                
                # 重置image_stream位置
                image_stream.seek(0)
                
                # 添加图片到幻灯片
                pic = slide.shapes.add_picture(
                    image_stream,
                    img_left, top,
                    width=Inches(actual_img_width),
                    height=Inches(actual_img_height)
                )
                
                # 添加图片边框
                line = pic.line
                line.color.rgb = (200, 200, 200)  # 浅灰色边框
                line.width = Pt(1)
                
                # 计算标签位置（图片下方）
                label_top = Inches(margin_top + actual_img_height + 0.1)
                
                # 添加项目名称标签（在图片下方）
                label_box = slide.shapes.add_textbox(
                    left, label_top,
                    Inches(img_width), Inches(label_height)
                )
                label_frame = label_box.text_frame
                label_frame.text = method['name']
                label_frame.paragraphs[0].font.size = Pt(14)
                label_frame.paragraphs[0].font.bold = True
                label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                label_frame.vertical_anchor = 1  # 顶部对齐
                
                logger.info(f"添加图片: {method['name']} - {metric} (尺寸: {actual_img_width:.2f}×{actual_img_height:.2f}英寸)")
                
            except Exception as e:
                logger.error(f"添加图片失败 {method['name']} - {metric}: {e}")
                continue
    
    def _add_comparison_layout_grid(self, slide, metric: str, methods: List[Dict]):
        """
        在幻灯片中添加网格布局的对比（适合更多项目）
        
        优化要点：
        1. 支持多行布局（自动计算最佳行列数）
        2. 项目名称标签放在图片下方
        3. 保持图片原始纵横比
        4. 添加图片边框以区分不同项目
        
        布局规则：
        - 2-3个项目：1行
        - 4-6个项目：2行
        - 7-9个项目：3行
        
        Args:
            slide: 幻灯片对象
            metric: 指标类型
            methods: 方法列表
        """
        n_methods = len(methods)
        
        # 计算最佳行列数
        if n_methods <= 3:
            n_rows = 1
            n_cols = n_methods
        elif n_methods <= 6:
            n_rows = 2
            n_cols = (n_methods + 1) // 2  # 向上取整
        else:
            n_rows = 3
            n_cols = (n_methods + 2) // 3  # 向上取整
        
        # 幻灯片尺寸
        slide_width = 10.0
        slide_height = 7.5
        
        # 留出空间
        margin_top = 1.2
        margin_bottom = 0.3
        margin_left = 0.3
        margin_right = 0.3
        
        # 可用空间
        available_width = slide_width - margin_left - margin_right
        available_height = slide_height - margin_top - margin_bottom
        
        # 标签高度
        label_height = 0.35
        
        # 间隙
        h_gap = 0.25  # 水平间隙
        v_gap = 0.3   # 垂直间隙
        
        # 计算每个单元格的尺寸
        cell_width = (available_width - h_gap * (n_cols - 1)) / n_cols
        cell_height = (available_height - v_gap * (n_rows - 1)) / n_rows
        
        # 图片最大尺寸（留出标签空间）
        max_img_width = cell_width
        max_img_height = cell_height - label_height - 0.1
        
        # 添加每个方法的热图
        for idx, method in enumerate(methods):
            # 计算当前图片的行列位置
            row = idx // n_cols
            col = idx % n_cols
            
            # 计算单元格左上角位置
            cell_left = margin_left + col * (cell_width + h_gap)
            cell_top = margin_top + row * (cell_height + v_gap)
            
            # 查找该方法的对应指标热图
            heatmap = next((hm for hm in method['heatmaps'] if hm['metric'] == metric), None)
            
            if not heatmap:
                # 如果该方法没有此指标，显示占位符
                placeholder = slide.shapes.add_textbox(
                    Inches(cell_left), Inches(cell_top),
                    Inches(cell_width), Inches(max_img_height)
                )
                text_frame = placeholder.text_frame
                text_frame.text = f"{method['name']}\n\n无此指标数据"
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame.paragraphs[0].font.size = Pt(12)
                text_frame.paragraphs[0].font.color.rgb = (128, 128, 128)
                continue
            
            # 添加图片
            try:
                # 将Base64解码为图片
                image_bytes = base64.b64decode(heatmap['image_data'])
                image_stream = io.BytesIO(image_bytes)
                
                # 获取图片原始尺寸以保持纵横比
                img = Image.open(io.BytesIO(image_bytes))
                orig_width, orig_height = img.size
                aspect_ratio = orig_height / orig_width
                
                # 计算实际图片尺寸（保持纵横比）
                actual_img_width = max_img_width
                actual_img_height = max_img_width * aspect_ratio
                
                # 如果高度超过最大高度，按高度缩放
                if actual_img_height > max_img_height:
                    actual_img_height = max_img_height
                    actual_img_width = max_img_height / aspect_ratio
                
                # 居中对齐
                img_left = cell_left + (cell_width - actual_img_width) / 2
                img_top = cell_top
                
                # 重置image_stream位置
                image_stream.seek(0)
                
                # 添加图片到幻灯片
                pic = slide.shapes.add_picture(
                    image_stream,
                    Inches(img_left), Inches(img_top),
                    width=Inches(actual_img_width),
                    height=Inches(actual_img_height)
                )
                
                # 添加图片边框
                line = pic.line
                line.color.rgb = (180, 180, 180)  # 浅灰色边框
                line.width = Pt(0.75)
                
                # 计算标签位置（图片下方，单元格内居中）
                label_top = cell_top + actual_img_height + 0.05
                
                # 添加项目名称标签（在图片下方）
                label_box = slide.shapes.add_textbox(
                    Inches(cell_left), Inches(label_top),
                    Inches(cell_width), Inches(label_height)
                )
                label_frame = label_box.text_frame
                label_frame.text = method['name']
                label_frame.paragraphs[0].font.size = Pt(12)
                label_frame.paragraphs[0].font.bold = True
                label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                label_frame.vertical_anchor = 1  # 顶部对齐
                
                logger.info(f"添加图片(网格): {method['name']} - {metric} (位置: 第{row+1}行第{col+1}列)")
                
            except Exception as e:
                logger.error(f"添加图片失败 {method['name']} - {metric}: {e}")
                continue
