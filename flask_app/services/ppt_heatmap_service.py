"""
PPT Image Replacement Service for the Immune Repertoire Analysis Web Application.
Provides functionality for replacing various analysis images in PowerPoint presentations.

Features:
1. Parse PPT template to identify image placeholders
2. Support multiple image types:
   - Sharing Analysis heatmaps (similarity heatmaps)
   - Network Plots (CDR3 sequence network graphs)
   - Isotype Upset Plots (isotype upset diagrams)
   - Tree Maps (tree diagrams)
3. Match generated images to PPT slide positions
4. Replace images while preserving layout
5. Auto-adjust layout based on sample count
6. Extract images from PPT for preview

Requirements: 10.4, 10.5
"""
import os
import io
import re
import base64
import logging
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass, field
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from flask_app.exceptions import (
    PPTError, PPTFileInvalidError, PPTParseError, 
    PPTSlideNotFoundError, PPTImageReplacementError, PPTNoHeatmapsError
)
from flask_app.services.image_border_manager import ImageBorderStyleManager

logger = logging.getLogger(__name__)


# Image type constants
IMAGE_TYPE_SHARING_ANALYSIS = 'sharing_analysis'
IMAGE_TYPE_NETWORK_PLOTS = 'network_plots'
IMAGE_TYPE_ISOTYPE_UPSET = 'isotype_upset'
IMAGE_TYPE_TREE_MAPS = 'tree_maps'


@dataclass
class SlideImageInfo:
    """Information about images in a slide"""
    slide_index: int
    slide_title: str
    image_type: str  # sharing_analysis, network_plots, isotype_upset, tree_maps
    chain_type: Optional[str] = None  # IGH, IGK, IGL, TRA, TRB, TRD, TRG (for sharing_analysis)
    metric_type: Optional[str] = None  # slide_type_1 (Expression/R²) or slide_type_2 (Morisita/Sorensen)
    slide_number_for_chain: int = 1  # 1 or 2 (which slide for this chain)
    sample_names: List[str] = field(default_factory=list)  # Sample names from table
    image_positions: List[Dict[str, Any]] = field(default_factory=list)
    image_data: List[Dict[str, Any]] = field(default_factory=list)  # Extracted image data


# Alias for backward compatibility
SlideHeatmapInfo = SlideImageInfo


@dataclass 
class ImageMapping:
    """Mapping between generated image and PPT position"""
    image_type: str  # sharing_analysis, network_plots, isotype_upset, tree_maps
    chain: Optional[str] = None  # For sharing_analysis
    metric: Optional[str] = None  # For sharing_analysis
    sample_name: Optional[str] = None  # For network_plots, isotype_upset
    slide_index: int = 0
    shape_index: int = 0
    position: Dict[str, float] = field(default_factory=dict)  # left, top, width, height in inches
    image_path: Optional[str] = None


# Alias for backward compatibility
HeatmapMapping = ImageMapping


@dataclass
class BorderConfig:
    """
    Configuration for image borders in PPT.
    
    This dataclass defines the border styling parameters that can be applied
    to images in PowerPoint presentations.
    
    Attributes:
        width_pt: Border width in points (pt). Default is 1.0pt.
        color_rgb: Border color as RGB tuple (r, g, b). Default is black (0, 0, 0).
                   Each value should be in range 0-255.
        dash_style: Border line style. Default is 'SOLID'.
    
    Requirements: 11.3, 11.4
    """
    width_pt: float = 1.0  # Border width in points
    color_rgb: Tuple[int, int, int] = (0, 0, 0)  # Border color RGB (black)
    dash_style: str = 'SOLID'  # Border line style
    
    def validate(self) -> bool:
        """
        Validate the border configuration.
        
        Checks that:
        - Width is positive and reasonable (0 < width <= 10pt)
        - RGB values are in valid range (0-255)
        - Dash style is a valid string
        
        Returns:
            True if configuration is valid, False otherwise
        
        Requirements: 11.3, 11.4
        """
        # Validate width - must be positive and reasonable
        if self.width_pt <= 0 or self.width_pt > 10:
            logger.error(f"Invalid border width: {self.width_pt}pt (must be 0 < width <= 10)")
            return False
        
        # Validate RGB values - must be in range 0-255
        r, g, b = self.color_rgb
        if not all(0 <= val <= 255 for val in [r, g, b]):
            logger.error(f"Invalid RGB values: {self.color_rgb} (must be 0-255)")
            return False
        
        # Validate dash style - must be a non-empty string
        if not isinstance(self.dash_style, str) or not self.dash_style:
            logger.error(f"Invalid dash style: {self.dash_style} (must be non-empty string)")
            return False
        
        return True


@dataclass
class ImageReplacementResult:
    """
    Result of image replacement operation in PPT.
    
    This dataclass encapsulates the results of replacing images in a PowerPoint
    presentation, including success status, counts, and any errors or warnings.
    
    Attributes:
        success: Whether the overall operation was successful
        replaced_count: Number of images successfully replaced
        total_count: Total number of images that were attempted to be replaced
        border_applied_count: Number of images that successfully received borders
        errors: List of error messages encountered during replacement
        warnings: List of warning messages encountered during replacement
    
    Requirements: 11.1, 11.2, 11.5
    """
    success: bool
    replaced_count: int
    total_count: int
    border_applied_count: int = 0  # Number of images with borders applied
    errors: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)
    
    def add_error(self, error: str) -> None:
        """Add an error message to the result."""
        self.errors.append(error)
        logger.error(f"Image replacement error: {error}")
    
    def add_warning(self, warning: str) -> None:
        """Add a warning message to the result."""
        self.warnings.append(warning)
        logger.warning(f"Image replacement warning: {warning}")
    
    def to_dict(self) -> Dict[str, Any]:
        """
        Convert result to dictionary for JSON serialization.
        
        Returns:
            Dictionary representation of the result
        """
        return {
            'success': self.success,
            'replaced_count': self.replaced_count,
            'total_count': self.total_count,
            'border_applied_count': self.border_applied_count,
            'errors': self.errors,
            'warnings': self.warnings,
            'replacement_rate': f"{self.replaced_count}/{self.total_count}",
            'border_rate': f"{self.border_applied_count}/{self.replaced_count}" if self.replaced_count > 0 else "0/0"
        }


class PPTImageService:
    """
    Service for replacing various analysis images in PowerPoint presentations.
    
    Supports the following image types:
    - Sharing Analysis heatmaps (7 chain types × 2 slides × 3 images)
    - Network Plots (sample-based network graphs)
    - Isotype Upset Plots (sample-based upset diagrams)
    - Tree Maps (tree diagrams)
    
    Supports the following chain types (for Sharing Analysis):
    - IGH, IGK, IGL (B-cell receptors)
    - TRA, TRB, TRD, TRG (T-cell receptors)
    """
    
    # Page title patterns for different image types
    PAGE_PATTERNS = {
        IMAGE_TYPE_NETWORK_PLOTS: [
            r'Network\s+Plots?\b',
            r'网络图',
        ],
        IMAGE_TYPE_ISOTYPE_UPSET: [
            r'Isotype\s+Upset\s+Plots?\b',
            r'同型体.*Upset',
        ],
        IMAGE_TYPE_TREE_MAPS: [
            r'Multi[-\s]?chain\s+Tree\s+Maps?\b',  # Must match "Multi-chain Tree Maps"
            r'树图',
        ],
    }
    
    # Configuration for skipping images per type
    # Network Plots: skip first image (usually a legend or overview)
    SKIP_FIRST_IMAGE = {
        IMAGE_TYPE_NETWORK_PLOTS: True,
        IMAGE_TYPE_ISOTYPE_UPSET: False,
        IMAGE_TYPE_TREE_MAPS: False,
        IMAGE_TYPE_SHARING_ANALYSIS: False,
    }
    
    # Chain type patterns for Sharing Analysis - Only match complete "Sharing Analysis - XXX" format
    CHAIN_PATTERNS = {
        'IGH': [r'Sharing\s+Analysis\s*[-–—]\s*IGH\b'],
        'IGK': [r'Sharing\s+Analysis\s*[-–—]\s*IGK\b'],
        'IGL': [r'Sharing\s+Analysis\s*[-–—]\s*IGL\b'],
        'TRA': [r'Sharing\s+Analysis\s*[-–—]\s*TRA\b'],
        'TRB': [r'Sharing\s+Analysis\s*[-–—]\s*TRB\b'],
        'TRD': [r'Sharing\s+Analysis\s*[-–—]\s*TRD\b'],
        'TRG': [r'Sharing\s+Analysis\s*[-–—]\s*TRG\b'],
    }
    
    # Metric type patterns for detection
    METRIC_PATTERNS = {
        'expression': [r'Expression', r'Reads'],
        'r2_outer': [r'R2.*Outer', r'R².*Outer', r'Outer.*R2', r'Outer'],
        'r2_inner': [r'R2.*Inner', r'R².*Inner', r'Inner.*R2', r'Inner'],
        'morisita_horn': [r'Morisita', r'Horn', r'MH'],
        'ucdr3': [r'uCDR3', r'CDR3', r'Unique'],
        'sorensen': [r'Sorensen', r'Sørensen', r'Dice'],
    }
    
    # Slide type definitions - each chain has 2 slides with 3 images each
    # Slide 1: Expression, R² Outer, R² Inner
    # Slide 2: Morisita-Horn, uCDR3, Sorensen
    SLIDE_TYPE_1_METRICS = ['expression', 'r2_outer', 'r2_inner']
    SLIDE_TYPE_2_METRICS = ['morisita_horn', 'ucdr3', 'sorensen']
    
    # Display names for metrics (Chinese)
    METRIC_DISPLAY_NAMES = {
        'expression': 'Expression',
        'r2_outer': 'R² Outer',
        'r2_inner': 'R² Inner',
        'morisita_horn': 'Morisita-Horn',
        'ucdr3': 'uCDR3',
        'sorensen': 'Sorensen',
    }
    
    def __init__(self):
        """Initialize the PPT image service"""
        self.presentation = None
        self.slide_image_info: List[SlideImageInfo] = []
        self.image_mappings: List[ImageMapping] = []
        self.layout_summary: Optional[Dict[str, Any]] = None
    
    # Backward compatibility aliases
    @property
    def slide_heatmap_info(self):
        return self.slide_image_info
    
    @slide_heatmap_info.setter
    def slide_heatmap_info(self, value):
        self.slide_image_info = value
    
    @property
    def heatmap_mappings(self):
        return self.image_mappings
    
    @heatmap_mappings.setter
    def heatmap_mappings(self, value):
        self.image_mappings = value
    
    def load_presentation(self, ppt_path: str) -> bool:
        """
        Load a PowerPoint presentation.
        
        Args:
            ppt_path: Path to the PPT file
            
        Returns:
            True if loaded successfully
            
        Raises:
            PPTFileInvalidError: If file does not exist or is invalid
            PPTParseError: If file cannot be parsed
        """
        if not os.path.exists(ppt_path):
            logger.error(f"PPT file not found: {ppt_path}")
            raise PPTFileInvalidError(
                message=f"PPT文件不存在: {ppt_path}",
                details={'path': ppt_path}
            )
        
        try:
            self.presentation = Presentation(ppt_path)
            self.ppt_path = ppt_path
            logger.info(f"Loaded presentation: {ppt_path} with {len(self.presentation.slides)} slides")
            return True
        except Exception as e:
            logger.error(f"Failed to load presentation: {e}")
            raise PPTParseError(
                message=f"无法解析PPT文件: {os.path.basename(ppt_path)}",
                details={'path': ppt_path, 'error': str(e)}
            )
    
    def load_presentation_from_bytes(self, ppt_bytes: bytes) -> bool:
        """
        Load a PowerPoint presentation from bytes.
        
        Args:
            ppt_bytes: PPT file content as bytes
            
        Returns:
            True if loaded successfully
            
        Raises:
            PPTFileInvalidError: If bytes are empty or invalid
            PPTParseError: If content cannot be parsed
        """
        if not ppt_bytes:
            logger.error("Empty PPT bytes provided")
            raise PPTFileInvalidError(
                message="PPT文件内容为空",
                details={'bytes_length': 0}
            )
        
        try:
            ppt_stream = io.BytesIO(ppt_bytes)
            self.presentation = Presentation(ppt_stream)
            self.ppt_path = None
            logger.info(f"Loaded presentation from bytes with {len(self.presentation.slides)} slides")
            return True
        except Exception as e:
            logger.error(f"Failed to load presentation from bytes: {e}")
            raise PPTParseError(
                message="无法解析PPT文件内容",
                details={'bytes_length': len(ppt_bytes), 'error': str(e)}
            )
    
    def analyze_slides(self) -> List[SlideImageInfo]:
        """
        Analyze all slides to identify image positions and types.
        
        Supports:
        1. Sharing Analysis slides (7 chains × 2 slides each)
        2. Network Plots slides (sample-based)
        3. Isotype Upset Plots slides (sample-based)
        4. Tree Maps slides
        
        Returns:
            List of SlideImageInfo objects
            
        Raises:
            PPTParseError: If no presentation is loaded or analysis fails
        """
        if not self.presentation:
            raise PPTParseError(
                message="未加载PPT文件，请先上传文件",
                details={'error': 'No presentation loaded'}
            )
        
        try:
            self.slide_image_info = []
            
            # Track slides by chain to determine slide type (1st or 2nd for each chain)
            chain_slide_count = {}
            
            for slide_idx, slide in enumerate(self.presentation.slides):
                slide_text = self._extract_slide_text(slide)
                
                # Check for different page types
                image_type = self._detect_page_type(slide_text)
                
                if image_type == IMAGE_TYPE_SHARING_ANALYSIS:
                    # Sharing Analysis page - detect chain type
                    chain_type = self._detect_chain_type(slide_text)
                    if not chain_type:
                        continue
                    
                    # Find image shapes
                    image_positions = self._find_image_shapes(slide)
                    if not image_positions:
                        continue
                    
                    # Track which slide number this is for the chain
                    if chain_type not in chain_slide_count:
                        chain_slide_count[chain_type] = 0
                    chain_slide_count[chain_type] += 1
                    slide_number_for_chain = chain_slide_count[chain_type]
                    
                    # Determine slide type
                    if slide_number_for_chain == 1:
                        slide_type = 'expression_r2'
                    elif slide_number_for_chain == 2:
                        slide_type = 'morisita_sorensen'
                    else:
                        detected_metrics = self._detect_metric_types(slide_text)
                        if any(m in ['r2_inner', 'r2_outer', 'expression_sharing'] for m in detected_metrics):
                            slide_type = 'expression_r2'
                        else:
                            slide_type = 'morisita_sorensen'
                    
                    info = SlideImageInfo(
                        slide_index=slide_idx,
                        slide_title=slide_text[:200],
                        image_type=IMAGE_TYPE_SHARING_ANALYSIS,
                        chain_type=chain_type,
                        metric_type=slide_type,
                        slide_number_for_chain=slide_number_for_chain,
                        image_positions=image_positions
                    )
                    self.slide_image_info.append(info)
                    logger.info(f"Found Sharing Analysis slide {slide_idx}: {chain_type} - {slide_type} ({slide_number_for_chain}/2) with {len(image_positions)} images")
                    
                elif image_type in [IMAGE_TYPE_NETWORK_PLOTS, IMAGE_TYPE_ISOTYPE_UPSET, IMAGE_TYPE_TREE_MAPS]:
                    # Sample-based pages - extract sample names from table
                    sample_names = self._extract_sample_names_from_table(slide)
                    image_positions = self._find_image_shapes(slide)
                    
                    if not image_positions:
                        continue
                    
                    # Skip first image for Network Plots (usually a legend or overview)
                    if self.SKIP_FIRST_IMAGE.get(image_type, False) and len(image_positions) > 1:
                        # Sort by position first to ensure we skip the correct one
                        sorted_positions = sorted(image_positions, key=lambda x: (x.get('top', 0), x.get('left', 0)))
                        image_positions = sorted_positions[1:]  # Skip the first image
                        logger.info(f"Skipped first image for {image_type}, remaining: {len(image_positions)}")
                    
                    info = SlideImageInfo(
                        slide_index=slide_idx,
                        slide_title=slide_text[:200],
                        image_type=image_type,
                        sample_names=sample_names,
                        image_positions=image_positions
                    )
                    self.slide_image_info.append(info)
                    logger.info(f"Found {image_type} slide {slide_idx} with {len(sample_names)} samples and {len(image_positions)} images")
            
            return self.slide_image_info
        except PPTParseError:
            raise
        except Exception as e:
            logger.error(f"Error analyzing slides: {e}")
            raise PPTParseError(
                message="分析PPT幻灯片时出错",
                details={'error': str(e)}
            )
    
    def _detect_page_type(self, text: str) -> Optional[str]:
        """Detect the page type from slide text"""
        # Check for Sharing Analysis first (most specific)
        for chain, patterns in self.CHAIN_PATTERNS.items():
            for pattern in patterns:
                if re.search(pattern, text, re.IGNORECASE):
                    return IMAGE_TYPE_SHARING_ANALYSIS
        
        # Check for other page types
        for page_type, patterns in self.PAGE_PATTERNS.items():
            for pattern in patterns:
                if re.search(pattern, text, re.IGNORECASE):
                    return page_type
        
        return None
    
    def _extract_sample_names_from_table(self, slide) -> List[str]:
        """Extract sample names from a table in the slide"""
        sample_names = []
        
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # Usually sample names are in the first row
                if table.rows:
                    first_row = table.rows[0]
                    for cell in first_row.cells:
                        cell_text = cell.text.strip()
                        if cell_text and not cell_text.lower() in ['sample', 'samples', '样本', '样本名']:
                            sample_names.append(cell_text)
        
        # If no table found, try to extract from text boxes
        if not sample_names:
            sample_names = self._extract_sample_names_from_text(slide)
        
        return sample_names
    
    def _extract_sample_names_from_text(self, slide) -> List[str]:
        """Extract sample names from text boxes in the slide (for images with labels)"""
        sample_names = []
        text_shapes = []
        
        # Collect all text shapes with their positions
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                # Skip common titles and headers
                skip_keywords = [
                    'network plot', 'isotype upset', 'tree map', 'sample', 
                    'sharing analysis', '网络图', '树图', 'plots', 'maps'
                ]
                if not any(kw.lower() in text.lower() for kw in skip_keywords):
                    # Check if this looks like a sample name (usually alphanumeric with underscores)
                    if self._looks_like_sample_name(text):
                        text_shapes.append({
                            'text': text,
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height
                        })
        
        # Sort by position (top to bottom, left to right)
        text_shapes.sort(key=lambda x: (x['top'], x['left']))
        
        for ts in text_shapes:
            if ts['text'] not in sample_names:
                sample_names.append(ts['text'])
        
        return sample_names
    
    def _looks_like_sample_name(self, text: str) -> bool:
        """Check if text looks like a sample name"""
        # Sample names typically:
        # - Contain letters and numbers or underscores
        # - Are not too long (usually < 50 chars)
        # - Don't contain too many spaces
        # - Often have patterns like XX_YY_ZZZZ
        
        if len(text) > 50 or len(text) < 2:
            return False
        
        # Count different character types
        has_letter = any(c.isalpha() for c in text)
        has_digit = any(c.isdigit() for c in text)
        space_count = text.count(' ')
        underscore_count = text.count('_')
        
        # Common sample name patterns
        if underscore_count >= 1 and has_letter:
            return True
        
        # Short alphanumeric strings
        if has_letter and has_digit and len(text) < 30 and space_count == 0:
            return True
        
        # Names with limited spaces
        if has_letter and space_count <= 1 and len(text) < 25:
            return True
        
        return False
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text from a slide"""
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
        return " ".join(text_parts)
    
    def _detect_chain_type(self, text: str) -> Optional[str]:
        """Detect chain type from slide text"""
        for chain, patterns in self.CHAIN_PATTERNS.items():
            for pattern in patterns:
                if re.search(pattern, text, re.IGNORECASE):
                    return chain
        return None
    
    def _detect_metric_types(self, text: str) -> List[str]:
        """Detect metric types from slide text"""
        detected = []
        for metric, patterns in self.METRIC_PATTERNS.items():
            for pattern in patterns:
                if re.search(pattern, text, re.IGNORECASE):
                    if metric not in detected:
                        detected.append(metric)
                    break
        return detected
    
    def _find_image_shapes(self, slide) -> List[Dict[str, Any]]:
        """Find all image shapes in a slide"""
        images = []
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append({
                    'shape_index': shape_idx,
                    'name': shape.name,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'left_inches': shape.left / Inches(1),
                    'top_inches': shape.top / Inches(1),
                    'width_inches': shape.width / Inches(1),
                    'height_inches': shape.height / Inches(1),
                })
        return images

    def extract_slide_images(self, slide_index: int, max_size: int = 200) -> List[Dict[str, Any]]:
        """
        Extract images from a specific slide as base64 encoded data.
        
        Args:
            slide_index: Index of the slide to extract images from
            max_size: Maximum dimension for thumbnail (default 200px)
            
        Returns:
            List of dictionaries with image data and metadata
            
        Raises:
            PPTParseError: If no presentation is loaded
            PPTSlideNotFoundError: If slide index is out of range
        """
        if not self.presentation:
            raise PPTParseError(
                message="未加载PPT文件",
                details={'error': 'No presentation loaded'}
            )
        
        if slide_index >= len(self.presentation.slides):
            raise PPTSlideNotFoundError(
                message=f"幻灯片索引超出范围: {slide_index}",
                details={
                    'slide_index': slide_index,
                    'total_slides': len(self.presentation.slides)
                }
            )
        
        slide = self.presentation.slides[slide_index]
        images = []
        
        # Find all picture shapes and sort by left position
        picture_shapes = []
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes.append((shape_idx, shape))
        
        # Sort by left position (left to right)
        picture_shapes.sort(key=lambda x: x[1].left)
        
        for shape_idx, shape in picture_shapes:
            try:
                # Get image blob from shape
                image_blob = shape.image.blob
                content_type = shape.image.content_type
                
                # Create thumbnail using PIL
                img = Image.open(io.BytesIO(image_blob))
                
                # Calculate thumbnail size maintaining aspect ratio
                width, height = img.size
                if width > height:
                    new_width = max_size
                    new_height = int(height * max_size / width)
                else:
                    new_height = max_size
                    new_width = int(width * max_size / height)
                
                # Create thumbnail
                img.thumbnail((new_width, new_height), Image.Resampling.LANCZOS)
                
                # Convert to base64
                buffer = io.BytesIO()
                img_format = 'PNG' if 'png' in content_type.lower() else 'JPEG'
                img.save(buffer, format=img_format)
                buffer.seek(0)
                base64_data = base64.b64encode(buffer.read()).decode('utf-8')
                
                # Determine mime type
                mime_type = 'image/png' if img_format == 'PNG' else 'image/jpeg'
                
                images.append({
                    'shape_index': shape_idx,
                    'name': shape.name,
                    'original_width': width,
                    'original_height': height,
                    'thumbnail_width': new_width,
                    'thumbnail_height': new_height,
                    'left_inches': shape.left / Inches(1),
                    'top_inches': shape.top / Inches(1),
                    'width_inches': shape.width / Inches(1),
                    'height_inches': shape.height / Inches(1),
                    'base64': base64_data,
                    'mime_type': mime_type,
                    'data_url': f"data:{mime_type};base64,{base64_data}"
                })
                
            except Exception as e:
                logger.error(f"Error extracting image from shape {shape_idx}: {e}")
                images.append({
                    'shape_index': shape_idx,
                    'name': shape.name,
                    'error': str(e),
                    'left_inches': shape.left / Inches(1),
                    'top_inches': shape.top / Inches(1),
                })
        
        return images

    def extract_all_heatmap_images(self, max_size: int = 200) -> None:
        """
        Extract images from all analyzed heatmap slides.
        
        Args:
            max_size: Maximum dimension for thumbnails
        """
        for slide_info in self.slide_heatmap_info:
            try:
                images = self.extract_slide_images(slide_info.slide_index, max_size)
                slide_info.image_data = images
            except Exception as e:
                logger.error(f"Error extracting images from slide {slide_info.slide_index}: {e}")
                slide_info.image_data = []

    def create_heatmap_mappings(
        self,
        available_heatmaps: Dict[str, Dict[str, str]]
    ) -> List[ImageMapping]:
        """
        Create mappings between available heatmaps and PPT positions.
        (Backward compatible method for Sharing Analysis heatmaps)
        
        Args:
            available_heatmaps: Dictionary mapping chain -> metric -> file_path
        
        Returns:
            List of ImageMapping objects
        """
        return self.create_image_mappings(
            sharing_analysis_images=available_heatmaps
        )
    
    def create_image_mappings(
        self,
        sharing_analysis_images: Optional[Dict[str, Dict[str, str]]] = None,
        network_plot_images: Optional[Dict[str, str]] = None,
        isotype_upset_images: Optional[Dict[str, str]] = None,
        tree_map_images: Optional[Dict[str, str]] = None,
        module: Optional[str] = None,
        layout_config: Optional[Dict[str, Any]] = None,
    ) -> List[ImageMapping]:
        """
        Create mappings between available images and PPT positions.
        
        Args:
            sharing_analysis_images: Dict mapping chain -> metric -> file_path
            network_plot_images: Dict mapping sample_name -> file_path
            isotype_upset_images: Dict mapping sample_name -> file_path
            tree_map_images: Dict mapping sample_name -> file_path
        
        Returns:
            List of ImageMapping objects
        """
        self.image_mappings = []
        self.layout_summary = None

        if module in [IMAGE_TYPE_NETWORK_PLOTS, IMAGE_TYPE_ISOTYPE_UPSET, IMAGE_TYPE_TREE_MAPS]:
            image_map = {
                IMAGE_TYPE_NETWORK_PLOTS: network_plot_images or {},
                IMAGE_TYPE_ISOTYPE_UPSET: isotype_upset_images or {},
                IMAGE_TYPE_TREE_MAPS: tree_map_images or {},
            }.get(module, {})
            module_slides = [
                slide_info for slide_info in self.slide_image_info
                if slide_info.image_type == module
            ]
            self.layout_summary = self._create_sample_based_mappings_with_layout(
                module_slides,
                image_map,
                module,
                layout_config or {}
            )
            return self.image_mappings
        
        for slide_info in self.slide_image_info:
            if module and slide_info.image_type != module:
                continue
            if slide_info.image_type == IMAGE_TYPE_SHARING_ANALYSIS:
                self._create_sharing_analysis_mappings(slide_info, sharing_analysis_images or {})
            elif slide_info.image_type == IMAGE_TYPE_NETWORK_PLOTS:
                self._create_sample_based_mappings(slide_info, network_plot_images or {}, IMAGE_TYPE_NETWORK_PLOTS, layout_config)
            elif slide_info.image_type == IMAGE_TYPE_ISOTYPE_UPSET:
                self._create_sample_based_mappings(slide_info, isotype_upset_images or {}, IMAGE_TYPE_ISOTYPE_UPSET, layout_config)
            elif slide_info.image_type == IMAGE_TYPE_TREE_MAPS:
                self._create_sample_based_mappings(slide_info, tree_map_images or {}, IMAGE_TYPE_TREE_MAPS, layout_config)
        
        return self.image_mappings
    
    def _create_sharing_analysis_mappings(
        self,
        slide_info: SlideImageInfo,
        available_heatmaps: Dict[str, Dict[str, str]]
    ):
        """Create mappings for Sharing Analysis heatmaps"""
        chain = slide_info.chain_type
        
        if chain not in available_heatmaps:
            logger.warning(f"No heatmaps available for chain {chain}")
            return
        
        chain_heatmaps = available_heatmaps[chain]
        
        # Match images to metrics based on position (left to right)
        sorted_images = sorted(
            slide_info.image_positions, 
            key=lambda x: x['left']
        )
        
        # Determine which metrics to map based on slide content
        slide_metrics = self._get_slide_metrics(slide_info)
        
        for img_idx, (image_info, metric) in enumerate(zip(sorted_images, slide_metrics)):
            if metric in chain_heatmaps:
                mapping = ImageMapping(
                    image_type=IMAGE_TYPE_SHARING_ANALYSIS,
                    chain=chain,
                    metric=metric,
                    slide_index=slide_info.slide_index,
                    shape_index=image_info['shape_index'],
                    position={
                        'left': image_info['left_inches'],
                        'top': image_info['top_inches'],
                        'width': image_info['width_inches'],
                        'height': image_info['height_inches'],
                    },
                    image_path=chain_heatmaps[metric]
                )
                self.image_mappings.append(mapping)
                logger.info(f"Mapped {chain}/{metric} to slide {slide_info.slide_index}, shape {image_info['shape_index']}")
    
    def _create_sample_based_mappings(
        self,
        slide_info: SlideImageInfo,
        available_images: Dict[str, str],
        image_type: str,
        layout_config: Optional[Dict[str, Any]] = None
    ):
        """Create mappings for sample-based images (Network Plots, Upset Plots, etc.)"""
        # Sort images by position (top to bottom, left to right)
        sorted_images = sorted(
            slide_info.image_positions,
            key=lambda x: (x['top'], x['left'])
        )
        
        # Match sample names to images
        sample_names = slide_info.sample_names
        
        for img_idx, image_info in enumerate(sorted_images):
            sample_name = sample_names[img_idx] if img_idx < len(sample_names) else None
            
            if sample_name and sample_name in available_images:
                mapping = ImageMapping(
                    image_type=image_type,
                    sample_name=sample_name,
                    slide_index=slide_info.slide_index,
                    shape_index=image_info['shape_index'],
                    position={
                        'left': image_info['left_inches'],
                        'top': image_info['top_inches'],
                        'width': image_info['width_inches'],
                        'height': image_info['height_inches'],
                    },
                    image_path=available_images[sample_name]
                )
                self.image_mappings.append(mapping)
                logger.info(f"Mapped {image_type}/{sample_name} to slide {slide_info.slide_index}, shape {image_info['shape_index']}")

    def _create_sample_based_mappings_with_layout(
        self,
        module_slides: List[SlideImageInfo],
        available_images: Dict[str, str],
        image_type: str,
        layout_config: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Create mappings for sample-based modules using an explicit layout configuration."""
        if not module_slides:
            return {
                'image_type': image_type,
                'total_samples': len(available_images),
                'mapped_samples': 0,
                'unmapped_samples': sorted(list(available_images.keys())),
                'pages_used': 0,
                'empty_slots': 0,
            }

        ordered_slide_infos = sorted(module_slides, key=lambda s: s.slide_index)
        positions_by_slide = []
        for slide_info in ordered_slide_infos:
            sorted_positions = sorted(
                slide_info.image_positions,
                key=lambda x: (x['top'], x['left'])
            )
            positions_by_slide.append((slide_info, sorted_positions))

        configured_order = layout_config.get('sample_order') or list(available_images.keys())
        ordered_samples = [sample for sample in configured_order if sample in available_images]
        for sample in available_images.keys():
            if sample not in ordered_samples:
                ordered_samples.append(sample)

        items_per_row = int(layout_config.get('items_per_row') or 3)
        items_per_row = max(1, min(items_per_row, 6))
        rows_per_page = int(layout_config.get('rows_per_page') or 2)
        rows_per_page = max(1, rows_per_page)
        configured_pages = layout_config.get('pages') or []

        mapped_samples = set()
        empty_slots = 0

        if configured_pages:
            for page_idx, page in enumerate(configured_pages):
                if page_idx >= len(positions_by_slide):
                    break
                slide_info, slide_positions = positions_by_slide[page_idx]
                slots = page.get('slots') or []
                for slot in slots:
                    sample_name = slot.get('sample_name')
                    if not sample_name or sample_name not in available_images:
                        continue
                    row = int(slot.get('row', 0))
                    col = int(slot.get('col', 0))
                    position_idx = row * items_per_row + col
                    if position_idx >= len(slide_positions):
                        continue
                    position = slide_positions[position_idx]
                    self.image_mappings.append(
                        ImageMapping(
                            image_type=image_type,
                            sample_name=sample_name,
                            slide_index=slide_info.slide_index,
                            shape_index=position['shape_index'],
                            position={
                                'left': position['left_inches'],
                                'top': position['top_inches'],
                                'width': position['width_inches'],
                                'height': position['height_inches'],
                            },
                            image_path=available_images[sample_name]
                        )
                    )
                    mapped_samples.add(sample_name)
        else:
            cursor = 0
            for slide_info, slide_positions in positions_by_slide:
                max_slots = min(len(slide_positions), items_per_row * rows_per_page)
                for pos_idx in range(max_slots):
                    if cursor >= len(ordered_samples):
                        empty_slots += (max_slots - pos_idx)
                        break
                    sample_name = ordered_samples[cursor]
                    cursor += 1
                    position = slide_positions[pos_idx]
                    self.image_mappings.append(
                        ImageMapping(
                            image_type=image_type,
                            sample_name=sample_name,
                            slide_index=slide_info.slide_index,
                            shape_index=position['shape_index'],
                            position={
                                'left': position['left_inches'],
                                'top': position['top_inches'],
                                'width': position['width_inches'],
                                'height': position['height_inches'],
                            },
                            image_path=available_images[sample_name]
                        )
                    )
                    mapped_samples.add(sample_name)

        unmapped_samples = [sample for sample in ordered_samples if sample not in mapped_samples]
        return {
            'image_type': image_type,
            'strategy': 'near_square',
            'items_per_row': items_per_row,
            'rows_per_page': rows_per_page,
            'total_samples': len(ordered_samples),
            'mapped_samples': len(mapped_samples),
            'unmapped_samples': unmapped_samples,
            'pages_used': len({m.slide_index for m in self.image_mappings if m.image_type == image_type}),
            'empty_slots': empty_slots,
        }
    
    def _get_slide_metrics(self, slide_info: SlideImageInfo) -> List[str]:
        """
        Determine the metrics for a slide based on its type.
        
        The PPT has two types of Sharing Analysis slides per chain:
        1. expression_r2 slides: expression, r2_outer, r2_inner
        2. morisita_sorensen slides: morisita_horn, ucdr3, sorensen
        """
        if slide_info.metric_type == 'expression_r2':
            return ['expression', 'r2_outer', 'r2_inner']
        elif slide_info.metric_type == 'morisita_sorensen':
            return ['morisita_horn', 'ucdr3', 'sorensen']
        
        # Fallback: detect from slide content
        slide_text = slide_info.slide_title.lower()
        
        # Check for Expression/R2 type slide
        if 'expression' in slide_text or 'r2' in slide_text.replace(' ', ''):
            return ['expression', 'r2_outer', 'r2_inner']
        
        # Check for Morisita/Sorensen type slide
        if 'morisita' in slide_text or 'sorensen' in slide_text or 'ucdr3' in slide_text:
            return ['morisita_horn', 'ucdr3', 'sorensen']
        
        # Default to expression type
        return ['expression', 'r2_outer', 'r2_inner']
    
    def replace_heatmaps(
        self,
        apply_borders: bool = True,
        border_config: Optional[BorderConfig] = None
    ) -> ImageReplacementResult:
        """
        Replace heatmap images in the presentation.
        
        Args:
            apply_borders: Whether to apply borders to replaced images (default: True)
            border_config: Optional border configuration
        
        Returns:
            ImageReplacementResult with replacement statistics
        """
        return self.replace_images(apply_borders=apply_borders, border_config=border_config)
    
    def replace_images(
        self, 
        apply_borders: bool = True,
        border_config: Optional[BorderConfig] = None
    ) -> ImageReplacementResult:
        """
        Replace all mapped images in the presentation.
        
        Args:
            apply_borders: Whether to apply borders to replaced images (default: True)
            border_config: Optional border configuration. If None, uses default black 1.0pt border.
        
        Returns:
            ImageReplacementResult with replacement statistics and border application counts
            
        Raises:
            PPTParseError: If no presentation is loaded
            PPTNoHeatmapsError: If no image mappings are created
            PPTImageReplacementError: If replacement fails
        
        Requirements: 11.1, 11.2, 11.5
        """
        if not self.presentation:
            raise PPTParseError(
                message="未加载PPT文件",
                details={'error': 'No presentation loaded'}
            )
        
        if not self.image_mappings:
            raise PPTNoHeatmapsError(
                message="未创建图片映射，请先调用 create_image_mappings",
                details={'error': 'No image mappings created'}
            )
        
        # Initialize border manager if borders should be applied
        border_manager = None
        if apply_borders:
            if border_config and border_config.validate():
                border_manager = ImageBorderStyleManager(
                    border_width=border_config.width_pt,
                    border_color=border_config.color_rgb
                )
            else:
                # Use default configuration (black, 1.0pt)
                border_manager = ImageBorderStyleManager()
            logger.info("Border manager initialized for image replacement")
        
        replaced_count = 0
        border_applied_count = 0
        total_count = len(self.image_mappings)
        errors = []
        warnings = []
        
        for mapping in self.image_mappings:
            image_path = mapping.image_path
            if not image_path or not os.path.exists(image_path):
                warning_msg = f"Image file not found: {image_path}"
                logger.warning(warning_msg)
                warnings.append(warning_msg)
                continue
            
            try:
                slide = self.presentation.slides[mapping.slide_index]
                
                # Find the shape to replace
                shape_to_replace = None
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        # Match by position (within tolerance)
                        if self._position_matches(shape, mapping.position):
                            shape_to_replace = shape
                            break
                
                if shape_to_replace:
                    # Get the position and size
                    left = shape_to_replace.left
                    top = shape_to_replace.top
                    width = shape_to_replace.width
                    height = shape_to_replace.height
                    
                    # Remove the old shape
                    sp = shape_to_replace._element
                    sp.getparent().remove(sp)
                    
                    # Add the new image
                    new_picture = slide.shapes.add_picture(
                        image_path,
                        left, top, width, height
                    )
                    
                    replaced_count += 1
                    
                    # Apply border to the newly added image
                    if border_manager:
                        try:
                            border_manager.apply_border(new_picture)
                            border_applied_count += 1
                            logger.debug(f"Applied border to replaced image on slide {mapping.slide_index}")
                        except Exception as border_error:
                            warning_msg = f"Failed to apply border on slide {mapping.slide_index}: {border_error}"
                            logger.warning(warning_msg)
                            warnings.append(warning_msg)
                    
                    # Log based on image type
                    if mapping.image_type == IMAGE_TYPE_SHARING_ANALYSIS:
                        logger.info(f"Replaced {mapping.chain}/{mapping.metric} on slide {mapping.slide_index}")
                    else:
                        logger.info(f"Replaced {mapping.image_type}/{mapping.sample_name} on slide {mapping.slide_index}")
                else:
                    warning_msg = f"Could not find shape to replace for mapping on slide {mapping.slide_index}"
                    logger.warning(warning_msg)
                    warnings.append(warning_msg)
                    
            except Exception as e:
                error_msg = f"Error replacing image on slide {mapping.slide_index}: {e}"
                logger.error(error_msg)
                errors.append(error_msg)
        
        # Create result object
        result = ImageReplacementResult(
            success=replaced_count > 0,
            replaced_count=replaced_count,
            total_count=total_count,
            border_applied_count=border_applied_count,
            errors=errors,
            warnings=warnings
        )
        
        logger.info(
            f"Image replacement complete: {replaced_count}/{total_count} replaced, "
            f"{border_applied_count} borders applied"
        )
        
        return result
    
    def _position_matches(self, shape, position: Dict[str, float], tolerance: float = 0.5) -> bool:
        """Check if shape position matches the expected position within tolerance (inches)"""
        shape_left = shape.left / Inches(1)
        shape_top = shape.top / Inches(1)
        
        return (
            abs(shape_left - position['left']) < tolerance and
            abs(shape_top - position['top']) < tolerance
        )
    
    def save_presentation(self, output_path: str) -> str:
        """
        Save the modified presentation.
        
        Args:
            output_path: Path to save the modified PPT
            
        Returns:
            Path to the saved file
            
        Raises:
            PPTParseError: If no presentation is loaded
            PPTError: If save operation fails
        """
        if not self.presentation:
            raise PPTParseError(
                message="未加载PPT文件",
                details={'error': 'No presentation loaded'}
            )
        
        try:
            # Ensure directory exists
            os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
            
            self.presentation.save(output_path)
            logger.info(f"Saved presentation to: {output_path}")
            return output_path
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            raise PPTError(
                message=f"保存PPT文件失败: {os.path.basename(output_path)}",
                details={'path': output_path, 'error': str(e)}
            )
    
    def save_presentation_to_bytes(self) -> bytes:
        """
        Save the modified presentation to bytes.
        
        Returns:
            PPT file content as bytes
            
        Raises:
            PPTParseError: If no presentation is loaded
            PPTError: If save operation fails
        """
        if not self.presentation:
            raise PPTParseError(
                message="未加载PPT文件",
                details={'error': 'No presentation loaded'}
            )
        
        try:
            output_stream = io.BytesIO()
            self.presentation.save(output_stream)
            output_stream.seek(0)
            return output_stream.read()
        except Exception as e:
            logger.error(f"Failed to save presentation to bytes: {e}")
            raise PPTError(
                message="保存PPT文件到内存失败",
                details={'error': str(e)}
            )
    
    def save(self, output_path: str) -> str:
        """
        Alias for save_presentation for convenience.
        
        Args:
            output_path: Path to save the modified PPT
            
        Returns:
            Path to the saved file
        """
        return self.save_presentation(output_path)
    
    def add_summary_slide(self, summary: Dict[str, Any]) -> None:
        """
        Add a summary slide at the end of the presentation.
        
        Args:
            summary: Dictionary containing replacement summary data
                - total_replaced: Total number of images replaced
                - modules: List of module status dicts
                - generated_at: Timestamp when summary was generated
        
        Requirements: Req 7 - Download Enhancement
        """
        if not self.presentation:
            raise PPTParseError(
                message="未加载PPT文件",
                details={'error': 'No presentation loaded'}
            )
        
        try:
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            
            # Add a blank slide at the end
            blank_layout = self.presentation.slide_layouts[6]  # Blank layout
            slide = self.presentation.slides.add_slide(blank_layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = "图片替换摘要"
            title_para.font.size = Pt(28)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(0, 51, 102)
            title_para.alignment = PP_ALIGN.CENTER
            
            # Add summary content
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3), Inches(9), Inches(5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            # Total replaced
            total_replaced = summary.get('total_replaced', 0)
            para = content_frame.paragraphs[0]
            para.text = f"总计替换图片: {total_replaced} 张"
            para.font.size = Pt(18)
            para.font.bold = True
            para.space_after = Pt(12)
            
            # Module details
            modules = summary.get('modules', [])
            for module in modules:
                para = content_frame.add_paragraph()
                module_name = module.get('name', 'Unknown')
                replaced = module.get('replaced', False)
                count = module.get('count', 0)
                
                status_icon = "✓" if replaced else "○"
                status_text = f"{count} 张" if replaced else "未替换"
                
                para.text = f"  {status_icon} {module_name}: {status_text}"
                para.font.size = Pt(14)
                para.font.color.rgb = RGBColor(0, 128, 0) if replaced else RGBColor(128, 128, 128)
                para.space_after = Pt(6)
            
            # Add timestamp
            generated_at = summary.get('generated_at', '')
            if generated_at:
                try:
                    from datetime import datetime
                    dt = datetime.fromisoformat(generated_at.replace('Z', '+00:00'))
                    formatted_time = dt.strftime('%Y-%m-%d %H:%M:%S')
                except:
                    formatted_time = generated_at
                
                para = content_frame.add_paragraph()
                para.text = ""
                para.space_before = Pt(20)
                
                para = content_frame.add_paragraph()
                para.text = f"生成时间: {formatted_time}"
                para.font.size = Pt(12)
                para.font.color.rgb = RGBColor(128, 128, 128)
                para.font.italic = True
            
            # Add footer note
            footer_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
            )
            footer_frame = footer_box.text_frame
            footer_para = footer_frame.paragraphs[0]
            footer_para.text = "此页面由PPT图片替换工具自动生成"
            footer_para.font.size = Pt(10)
            footer_para.font.color.rgb = RGBColor(150, 150, 150)
            footer_para.alignment = PP_ALIGN.CENTER
            
            logger.info("Added summary slide to presentation")
            
        except Exception as e:
            logger.error(f"Failed to add summary slide: {e}")
            raise PPTError(
                message="添加摘要页失败",
                details={'error': str(e)}
            )
    
    def get_slide_summary(self, include_images: bool = False) -> List[Dict[str, Any]]:
        """
        Get a summary of analyzed slides.
        
        Args:
            include_images: Whether to include extracted image data (base64)
        
        Returns:
            List of slide summaries with detailed image information
        """
        # Get expected metrics for each slide type - using correct metric names
        slide_type_metrics = {
            'expression_r2': ['expression', 'r2_outer', 'r2_inner'],
            'morisita_sorensen': ['morisita_horn', 'ucdr3', 'sorensen']
        }
        
        # Image type display names
        image_type_display = {
            IMAGE_TYPE_SHARING_ANALYSIS: 'Sharing Analysis 热图',
            IMAGE_TYPE_NETWORK_PLOTS: 'Network Plots 网络图',
            IMAGE_TYPE_ISOTYPE_UPSET: 'Isotype Upset Plots',
            IMAGE_TYPE_TREE_MAPS: 'Tree Maps 树图',
        }
        
        summaries = []
        for info in self.slide_image_info:
            # IMPORTANT: Sort by left position only for Sharing Analysis (horizontal layout)
            # This ensures Expression -> R² Outer -> R² Inner order (left to right)
            if info.image_type == IMAGE_TYPE_SHARING_ANALYSIS:
                sorted_images = sorted(info.image_positions, key=lambda x: x.get('left', 0))
            else:
                # For sample-based pages, sort by top then left (grid layout)
                sorted_images = sorted(info.image_positions, key=lambda x: (x.get('top', 0), x.get('left', 0)))
            
            # Build image positions with assignments
            image_positions = []
            
            if info.image_type == IMAGE_TYPE_SHARING_ANALYSIS:
                # Sharing Analysis - assign metrics
                expected_metrics = slide_type_metrics.get(info.metric_type, [])
                for idx, img in enumerate(sorted_images):
                    metric = expected_metrics[idx] if idx < len(expected_metrics) else 'unknown'
                    metric_display = self.METRIC_DISPLAY_NAMES.get(metric, metric)
                    
                    img_info = self._build_image_info(img, idx, include_images, info)
                    img_info['metric'] = metric
                    img_info['metric_display'] = metric_display
                    image_positions.append(img_info)
                
                summary = {
                    'slide_index': info.slide_index,
                    'image_type': info.image_type,
                    'image_type_display': image_type_display.get(info.image_type, info.image_type),
                    'chain_type': info.chain_type,
                    'metric_type': info.metric_type,
                    'slide_number_for_chain': info.slide_number_for_chain,
                    'slide_type_display': 'Expression/R² 类型' if info.metric_type == 'expression_r2' else 'Morisita/Sorensen 类型',
                    'expected_metrics': expected_metrics,
                    'expected_metrics_display': [self.METRIC_DISPLAY_NAMES.get(m, m) for m in expected_metrics],
                    'image_count': len(info.image_positions),
                    'image_positions': image_positions,
                }
            else:
                # Sample-based pages - assign sample names
                for idx, img in enumerate(sorted_images):
                    sample_name = info.sample_names[idx] if idx < len(info.sample_names) else f'Sample {idx + 1}'
                    
                    img_info = self._build_image_info(img, idx, include_images, info)
                    img_info['sample_name'] = sample_name
                    image_positions.append(img_info)
                
                summary = {
                    'slide_index': info.slide_index,
                    'image_type': info.image_type,
                    'image_type_display': image_type_display.get(info.image_type, info.image_type),
                    'sample_names': info.sample_names,
                    'sample_count': len(info.sample_names),
                    'image_count': len(info.image_positions),
                    'image_positions': image_positions,
                }
            
            summaries.append(summary)
        
        return summaries
    
    def _build_image_info(self, img: Dict, idx: int, include_images: bool, slide_info: SlideImageInfo) -> Dict[str, Any]:
        """Build image info dictionary"""
        img_info = {
            'index': idx,
            'shape_index': img['shape_index'],
            'name': img.get('name', ''),
            'position': f"({img['left_inches']:.1f}, {img['top_inches']:.1f})",
            'size': f"{img['width_inches']:.1f} x {img['height_inches']:.1f}",
            'left_inches': img['left_inches'],
            'top_inches': img['top_inches'],
            'width_inches': img['width_inches'],
            'height_inches': img['height_inches'],
        }
        
        # Add image data if available and requested
        if include_images and slide_info.image_data:
            logger.debug(f"Looking for shape_index={img['shape_index']} in {len(slide_info.image_data)} image_data items")
            for img_data in slide_info.image_data:
                if img_data.get('shape_index') == img['shape_index']:
                    img_info['data_url'] = img_data.get('data_url')
                    img_info['thumbnail_width'] = img_data.get('thumbnail_width')
                    img_info['thumbnail_height'] = img_data.get('thumbnail_height')
                    logger.debug(f"Found matching image data for shape_index={img['shape_index']}, data_url length: {len(img_info.get('data_url', ''))}")
                    break
            else:
                logger.warning(f"No matching image data found for shape_index={img['shape_index']}")
        elif include_images:
            logger.warning(f"include_images=True but slide_info.image_data is empty for slide {slide_info.slide_index}")
        
        return img_info
    
    def get_mapping_summary(self) -> List[Dict[str, Any]]:
        """
        Get a summary of image mappings.
        
        Returns:
            List of mapping summaries
        """
        summaries = []
        for m in self.image_mappings:
            summary = {
                'image_type': m.image_type,
                'slide_index': m.slide_index,
                'shape_index': m.shape_index,
                'image_file': os.path.basename(m.image_path) if m.image_path else None,
                'has_file': m.image_path is not None and os.path.exists(m.image_path) if m.image_path else False
            }
            
            if m.image_type == IMAGE_TYPE_SHARING_ANALYSIS:
                summary['chain'] = m.chain
                summary['metric'] = m.metric
            else:
                summary['sample_name'] = m.sample_name
            
            summaries.append(summary)
        
        return summaries

    def render_slide_preview(self, slide_index: int, max_size: int = 800) -> Dict[str, Any]:
        """
        Render a slide as an image for preview.
        
        Note: python-pptx doesn't support direct slide rendering.
        This method extracts slide information and images for preview.
        
        Args:
            slide_index: Index of the slide to render
            max_size: Maximum dimension for the preview
            
        Returns:
            Dictionary with slide preview information
        """
        if not self.presentation:
            raise PPTParseError(
                message="未加载PPT文件",
                details={'error': 'No presentation loaded'}
            )
        
        if slide_index >= len(self.presentation.slides):
            raise PPTSlideNotFoundError(
                message=f"幻灯片索引超出范围: {slide_index}",
                details={
                    'slide_index': slide_index,
                    'total_slides': len(self.presentation.slides)
                }
            )
        
        slide = self.presentation.slides[slide_index]
        
        # Extract slide title
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                title = shape.text.strip()[:100]
                break
        
        # Count images
        image_count = 0
        images_data = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                try:
                    # Extract image thumbnail
                    image_blob = shape.image.blob
                    content_type = shape.image.content_type
                    
                    img = Image.open(io.BytesIO(image_blob))
                    
                    # Create thumbnail
                    thumb_size = min(150, max_size // 4)
                    img.thumbnail((thumb_size, thumb_size), Image.Resampling.LANCZOS)
                    
                    buffer = io.BytesIO()
                    img_format = 'PNG' if 'png' in content_type.lower() else 'JPEG'
                    img.save(buffer, format=img_format)
                    buffer.seek(0)
                    base64_data = base64.b64encode(buffer.read()).decode('utf-8')
                    
                    mime_type = 'image/png' if img_format == 'PNG' else 'image/jpeg'
                    
                    images_data.append({
                        'shape_index': shape_idx,
                        'data_url': f"data:{mime_type};base64,{base64_data}",
                        'position': f"({shape.left / Inches(1):.1f}, {shape.top / Inches(1):.1f})"
                    })
                except Exception as e:
                    logger.error(f"Error extracting image from slide {slide_index}, shape {shape_idx}: {e}")
        
        # Create a composite preview image if there are images
        preview_data_url = None
        if images_data:
            # Use the first image as the preview (or create a composite)
            preview_data_url = images_data[0]['data_url'] if images_data else None
        
        return {
            'index': slide_index,
            'title': title or f'Slide {slide_index + 1}',
            'has_images': image_count > 0,
            'image_count': image_count,
            'images': images_data,
            'data_url': preview_data_url
        }


# Backward compatibility alias
PPTHeatmapService = PPTImageService


def replace_ppt_heatmaps(
    ppt_path: str,
    heatmap_dir: str,
    output_path: Optional[str] = None
) -> Tuple[str, int, List[Dict]]:
    """
    Convenience function to replace heatmaps in a PPT file.
    (Backward compatible function)
    
    Args:
        ppt_path: Path to the input PPT file
        heatmap_dir: Directory containing generated heatmaps
        output_path: Path for output file (default: adds '_updated' suffix)
        
    Returns:
        Tuple of (output_path, replaced_count, mapping_summary)
    """
    return replace_ppt_images(ppt_path, heatmap_dir, output_path)


def replace_ppt_images(
    ppt_path: str,
    image_dir: str,
    output_path: Optional[str] = None,
    image_types: Optional[List[str]] = None
) -> Tuple[str, int, List[Dict]]:
    """
    Convenience function to replace images in a PPT file.
    
    Args:
        ppt_path: Path to the input PPT file
        image_dir: Directory containing generated images
        output_path: Path for output file (default: adds '_updated' suffix)
        image_types: List of image types to replace (default: all)
        
    Returns:
        Tuple of (output_path, replaced_count, mapping_summary)
    """
    service = PPTImageService()
    
    # Load presentation
    service.load_presentation(ppt_path)
    
    # Analyze slides
    service.analyze_slides()
    
    # Scan image directory and build available images dict
    sharing_analysis_images = scan_heatmap_directory(image_dir)
    network_plot_images = scan_sample_images(image_dir, 'network_plots')
    isotype_upset_images = scan_sample_images(image_dir, 'isotype_upset')
    tree_map_images = scan_sample_images(image_dir, 'tree_maps')
    
    # Create mappings
    service.create_image_mappings(
        sharing_analysis_images=sharing_analysis_images,
        network_plot_images=network_plot_images,
        isotype_upset_images=isotype_upset_images,
        tree_map_images=tree_map_images
    )
    
    # Replace images
    result = service.replace_images()
    
    # Save
    if output_path is None:
        base, ext = os.path.splitext(ppt_path)
        output_path = f"{base}_updated{ext}"
    
    service.save_presentation(output_path)
    
    # Log border application details
    logger.info(
        f"Image replacement complete: {result.replaced_count}/{result.total_count} replaced, "
        f"{result.border_applied_count} borders applied"
    )
    
    return output_path, result.replaced_count, service.get_mapping_summary()


def scan_heatmap_directory(heatmap_dir: str) -> Dict[str, Dict[str, str]]:
    """
    Scan a directory for heatmap files and organize by chain and metric.
    
    Expected file naming convention:
    - {chain}_{metric}_heatmap.png (e.g., IGH_r2_inner_heatmap.png)
    - {chain}/{metric}_heatmap.png (e.g., IGH/r2_inner_heatmap.png)
    - {chain}/{metric}_sharing_heatmap.png (e.g., IGH/expression_sharing_heatmap.png)
    - {chain}_{metric}.png (e.g., TRB_r2_inner.png) - from similarity_heatmap module
    
    Args:
        heatmap_dir: Directory to scan
        
    Returns:
        Dictionary mapping chain -> metric -> file_path
    """
    available = {}
    
    chains = ['IGH', 'IGK', 'IGL', 'TRA', 'TRB', 'TRD', 'TRG']
    
    # Metric name mappings: internal_name -> [possible file patterns]
    # This handles different naming conventions in the generated files
    # Including patterns from similarity_heatmap module output
    # Priority order: most specific first
    metric_patterns = {
        'expression': [
            'expression_sharing_heatmap.png',  # Primary format from analysis
            'expression_heatmap.png', 
            'expression_sharing.png',
            'expression.png'
        ],
        'r2_inner': [
            'r2_inner_heatmap.png',  # Primary format from analysis
            'r2_inner_sharing_heatmap.png',
            'r2_inner.png'
        ],
        'r2_outer': [
            'r2_outer_heatmap.png',  # Primary format from analysis
            'r2_outer_sharing_heatmap.png',
            'r2_outer.png'
        ],
        'morisita_horn': [
            'morisita_horn_heatmap.png',  # Primary format from analysis
            'morisita_heatmap.png',
            'morisita_horn.png', 
            'morisita.png'
        ],
        'ucdr3': [
            'cdr3_sharing_heatmap.png',  # Primary format from analysis
            'ucdr3_heatmap.png', 
            'cdr3_heatmap.png',
            'ucdr3.png',
            'cdr3_sharing.png'
        ],
        'sorensen': [
            'sorensen_heatmap.png',  # Primary format from analysis
            'sorensen.png'
        ],
    }
    
    heatmap_path = Path(heatmap_dir)
    
    # Check if directory exists
    if not heatmap_path.exists():
        logger.warning(f"Heatmap directory does not exist: {heatmap_dir}")
        return available
    
    for chain in chains:
        chain_heatmaps = {}
        
        # Check for chain subdirectory (e.g., analysis_xxx/TRB/ or original_results/IGH/)
        chain_dir = heatmap_path / chain
        if chain_dir.exists() and chain_dir.is_dir():
            # Check if there's a heatmaps subdirectory (new format from downloadAllAsZip)
            heatmaps_subdir = chain_dir / 'heatmaps'
            search_dirs = [heatmaps_subdir] if heatmaps_subdir.exists() else [chain_dir]
            
            for search_dir in search_dirs:
                # First, try direct file matching with expected patterns
                for metric, patterns in metric_patterns.items():
                    if metric in chain_heatmaps:
                        continue
                    for pattern in patterns:
                        file_path = search_dir / pattern
                        if file_path.exists():
                            chain_heatmaps[metric] = str(file_path)
                            logger.debug(f"Found heatmap (direct): {chain}/{metric} -> {file_path}")
                            break
                        # Also try with chain prefix (e.g., IGH_expression_sharing_heatmap.png)
                        file_path = search_dir / f"{chain}_{pattern}"
                        if file_path.exists():
                            chain_heatmaps[metric] = str(file_path)
                            logger.debug(f"Found heatmap (direct with prefix): {chain}/{metric} -> {file_path}")
                            break
                
                # If not all metrics found, scan directory for remaining
                if len(chain_heatmaps) < len(metric_patterns):
                    for file_path in search_dir.glob('*.png'):
                        filename = file_path.name.lower()
                        for metric, patterns in metric_patterns.items():
                            if metric in chain_heatmaps:
                                continue
                            for pattern in patterns:
                                if filename == pattern.lower() or filename == f"{chain.lower()}_{pattern}".lower():
                                    chain_heatmaps[metric] = str(file_path)
                                    logger.debug(f"Found heatmap (scan): {chain}/{metric} -> {file_path}")
                                    break
        
        # Also check root directory with chain prefix
        for metric, patterns in metric_patterns.items():
            if metric not in chain_heatmaps:
                for pattern in patterns:
                    # Try with chain prefix in root directory
                    file_path = heatmap_path / f"{chain}_{pattern}"
                    if file_path.exists():
                        chain_heatmaps[metric] = str(file_path)
                        logger.debug(f"Found heatmap: {chain}/{metric} -> {file_path}")
                        break
                    # Try lowercase chain
                    file_path = heatmap_path / f"{chain.lower()}_{pattern}"
                    if file_path.exists():
                        chain_heatmaps[metric] = str(file_path)
                        logger.debug(f"Found heatmap: {chain}/{metric} -> {file_path}")
                        break
        
        # Final fallback: scan root directory for any files matching chain_metric pattern
        if not chain_heatmaps:
            for file_path in heatmap_path.glob(f'{chain}_*.png'):
                filename = file_path.stem.lower()
                for metric in metric_patterns.keys():
                    if metric in chain_heatmaps:
                        continue
                    metric_lower = metric.replace('_', '')
                    if metric in filename or metric_lower in filename:
                        chain_heatmaps[metric] = str(file_path)
                        logger.debug(f"Found heatmap (fallback): {chain}/{metric} -> {file_path}")
            
            # Also try lowercase chain
            for file_path in heatmap_path.glob(f'{chain.lower()}_*.png'):
                filename = file_path.stem.lower()
                for metric in metric_patterns.keys():
                    if metric in chain_heatmaps:
                        continue
                    metric_lower = metric.replace('_', '')
                    if metric in filename or metric_lower in filename:
                        chain_heatmaps[metric] = str(file_path)
                        logger.debug(f"Found heatmap (fallback): {chain}/{metric} -> {file_path}")
        
        if chain_heatmaps:
            available[chain] = chain_heatmaps
            logger.info(f"Found {len(chain_heatmaps)} heatmaps for chain {chain}")
    
    if not available:
        logger.warning(f"No heatmaps found in directory: {heatmap_dir}")
    else:
        total = sum(len(v) for v in available.values())
        logger.info(f"Total heatmaps found: {total} across {len(available)} chains")
    
    return available


def scan_sample_images(image_dir: str, image_type: str) -> Dict[str, str]:
    """
    Scan a directory for sample-based images.
    
    Expected directory structures:
    1. Flat structure: {image_dir}/{sample_name}_treemap.png
    2. Subdirectory structure: {image_dir}/individual_treemaps/{sample_name}_treemap.png
    3. Sample folder structure (extracted_images): {image_dir}/{sample_name}/image_-1.png (network)
       - image_-1.png = network plot
       - image_0.png, image_1.png, etc. = isotype upset plots
    
    Args:
        image_dir: Base directory to scan
        image_type: Type of images to scan for (network_plots, isotype_upset, tree_maps)
        
    Returns:
        Dictionary mapping sample_name -> file_path
    """
    available = {}
    
    # Map image type to directory names, file prefixes, suffixes and patterns
    type_config = {
        'network_plots': {
            'dirs': ['network_plots', 'network', 'networks', 'network_plot', ''],
            'prefixes': ['network_', 'net_', 'cdr3_network_'],
            'suffixes': ['_network', '_net', '_cdr3_network'],
            'patterns': ['*_cdr3_network.png', '*_cdr3_network.jpg', '*_network.png'],
            'sample_folder_file': 'image_-1.png'  # In extracted_images format: image_-1.png = network
        },
        'isotype_upset': {
            'dirs': ['isotype_upset', 'upset', 'isotype_upset_plots', 'upset_plots', ''],
            'prefixes': ['upset_', 'isotype_', 'isotype_upset_'],
            'suffixes': ['_upset', '_isotype', '_isotype_upset'],
            'patterns': ['*_isotype_upset.png', '*_isotype_upset.jpg', '*_upset.png'],
            'sample_folder_files': ['image_0.png', 'image_1.png', 'image_2.png']  # In extracted_images format
        },
        'tree_maps': {
            'dirs': ['tree_maps', 'treemaps', 'tree', 'tree_map', 'individual_treemaps', ''],
            'prefixes': ['tree_', 'treemap_', 'tree_map_'],
            'suffixes': ['_tree', '_treemap', '_tree_map'],
            'patterns': ['*_treemap.png', '*_treemap.jpg', '*.png']  # Also match any .png in treemap dirs
        },
    }
    
    config = type_config.get(image_type, {'dirs': [image_type], 'prefixes': [], 'suffixes': [], 'patterns': []})
    base_path = Path(image_dir)
    
    def extract_sample_name(filename: str, prefixes: List[str], suffixes: List[str]) -> str:
        """Extract sample name from filename by removing known prefixes/suffixes."""
        name = filename
        
        # Remove prefixes
        for prefix in prefixes:
            if name.lower().startswith(prefix.lower()):
                name = name[len(prefix):]
                break
        
        # Remove suffixes
        for suffix in suffixes:
            if name.lower().endswith(suffix.lower()):
                name = name[:-len(suffix)]
                break
        
        return name
    
    # Check if directory exists
    if not base_path.exists():
        logger.warning(f"Directory does not exist: {image_dir}")
        return available
    
    # Strategy 1: Check for sample folder structure (extracted_images format)
    # Each sample has its own folder with image_-1.png (network), image_0.png (upset), etc.
    sample_folders_found = False
    for item in base_path.iterdir():
        if item.is_dir():
            # Check if this looks like a sample folder
            if image_type == 'network_plots':
                # Look for image_-1.png in sample folder
                network_file = item / 'image_-1.png'
                if network_file.exists():
                    available[item.name] = str(network_file)
                    sample_folders_found = True
            elif image_type == 'isotype_upset':
                # Look for image_0.png, image_1.png, etc. in sample folder
                for upset_file_name in config.get('sample_folder_files', []):
                    upset_file = item / upset_file_name
                    if upset_file.exists():
                        # Use sample_name + index as key to support multiple upset images per sample
                        idx = upset_file_name.replace('image_', '').replace('.png', '')
                        sample_key = f"{item.name}_upset_{idx}" if idx != '0' else item.name
                        available[sample_key] = str(upset_file)
                        sample_folders_found = True
    
    if sample_folders_found:
        logger.info(f"Found {len(available)} {image_type} images using sample folder structure from {image_dir}")
        return available
    
    # Strategy 2: Try different directory names for flat/subdirectory structure
    found_images = False
    for dir_name in config['dirs']:
        if dir_name == '':
            type_dir = base_path
        else:
            type_dir = base_path / dir_name
        
        if not type_dir.exists() or not type_dir.is_dir():
            continue
        
        # First try pattern-based matching (e.g., *_cdr3_network.png, *_treemap.png)
        for pattern in config.get('patterns', []):
            for file_path in type_dir.glob(pattern):
                sample_name = extract_sample_name(
                    file_path.stem, 
                    config['prefixes'], 
                    config['suffixes']
                )
                if sample_name and sample_name not in available:
                    available[sample_name] = str(file_path)
                    found_images = True
        
        # If no pattern matches, scan all image files in the directory
        if not found_images and dir_name != '':
            for file_path in type_dir.glob('*.png'):
                sample_name = extract_sample_name(
                    file_path.stem, 
                    config['prefixes'], 
                    config['suffixes']
                )
                if sample_name and sample_name not in available:
                    available[sample_name] = str(file_path)
                    found_images = True
            
            for file_path in type_dir.glob('*.jpg'):
                sample_name = extract_sample_name(
                    file_path.stem,
                    config['prefixes'],
                    config['suffixes']
                )
                if sample_name and sample_name not in available:
                    available[sample_name] = str(file_path)
                    found_images = True
        
        if found_images:
            break  # Found images, stop searching
    
    # Strategy 3: If still no images found, scan root directory with pattern-based matching
    if not available:
        for pattern in config.get('patterns', []):
            for file_path in base_path.glob(pattern):
                sample_name = extract_sample_name(
                    file_path.stem,
                    config['prefixes'],
                    config['suffixes']
                )
                if sample_name and sample_name not in available:
                    available[sample_name] = str(file_path)
    
    logger.info(f"Scanned {image_type} images from {image_dir}: found {len(available)} samples")
    return available


def calculate_auto_layout(sample_count: int, page_width: float = 13.33, page_height: float = 7.5, single_page: bool = True) -> Dict[str, Any]:
    """
    Calculate optimal layout for sample-based images.
    
    For Network Plots, Isotype Upset Plots, and Tree Maps:
    - All images are placed on a single page with auto-adjustment
    - Layout dynamically adjusts based on sample count
    
    Args:
        sample_count: Number of samples to display
        page_width: Page width in inches (default: 13.33 for widescreen)
        page_height: Page height in inches (default: 7.5 for widescreen)
        single_page: If True, fit all images on single page (default: True)
        
    Returns:
        Dictionary with layout information:
        - rows: number of rows
        - cols: number of columns per row
        - image_width: width of each image
        - image_height: height of each image
        - positions: list of (left, top) positions for each image
    """
    import math
    
    # Margins and spacing
    margin_left = 0.4
    margin_right = 0.4
    margin_top = 0.8  # Space for title
    margin_bottom = 0.4
    spacing_h = 0.15  # Horizontal spacing between images
    spacing_v = 0.2  # Vertical spacing between images
    
    available_width = page_width - margin_left - margin_right
    available_height = page_height - margin_top - margin_bottom
    
    if sample_count <= 0:
        return {
            'rows': 0,
            'cols': 0,
            'image_width': 0,
            'image_height': 0,
            'positions': [],
            'samples_per_page': 0,
            'total_pages': 0
        }
    
    if single_page:
        # Calculate optimal grid for fitting all samples on one page
        # Try different row/column combinations and pick the best one
        best_layout = None
        best_area = 0
        
        for rows in range(1, sample_count + 1):
            cols = math.ceil(sample_count / rows)
            
            # Calculate image size for this layout
            total_h_spacing = spacing_h * (cols - 1)
            total_v_spacing = spacing_v * (rows - 1)
            
            img_width = (available_width - total_h_spacing) / cols
            img_height = (available_height - total_v_spacing) / rows
            
            # Keep aspect ratio reasonable (between 0.5 and 2.0)
            aspect_ratio = img_width / img_height if img_height > 0 else 0
            if aspect_ratio < 0.3 or aspect_ratio > 3.0:
                continue
            
            # Calculate area utilized
            area = img_width * img_height
            
            if area > best_area:
                best_area = area
                best_layout = {
                    'rows': rows,
                    'cols': cols,
                    'image_width': img_width,
                    'image_height': img_height
                }
        
        if best_layout:
            rows = best_layout['rows']
            cols = best_layout['cols']
            image_width = best_layout['image_width']
            image_height = best_layout['image_height']
        else:
            # Fallback: simple grid
            cols = math.ceil(math.sqrt(sample_count))
            rows = math.ceil(sample_count / cols)
            total_h_spacing = spacing_h * (cols - 1)
            total_v_spacing = spacing_v * (rows - 1)
            image_width = (available_width - total_h_spacing) / cols
            image_height = (available_height - total_v_spacing) / rows
    else:
        # Original multi-page logic
        if sample_count <= 3:
            rows, cols = 1, sample_count
        elif sample_count <= 6:
            rows = 2
            cols = (sample_count + 1) // 2
        elif sample_count <= 12:
            rows = 3
            cols = 4
        else:
            rows = 3
            cols = 4
        
        total_h_spacing = spacing_h * (cols - 1)
        total_v_spacing = spacing_v * (rows - 1)
        image_width = (available_width - total_h_spacing) / cols
        image_height = (available_height - total_v_spacing) / rows
    
    # Calculate positions
    positions = []
    for row in range(rows):
        for col in range(cols):
            idx = row * cols + col
            if idx >= sample_count:
                break
            
            left = margin_left + col * (image_width + spacing_h)
            top = margin_top + row * (image_height + spacing_v)
            positions.append({
                'left': left,
                'top': top,
                'width': image_width,
                'height': image_height,
                'index': idx
            })
    
    return {
        'rows': rows,
        'cols': cols,
        'image_width': image_width,
        'image_height': image_height,
        'positions': positions,
        'samples_per_page': rows * cols,
        'total_pages': 1 if single_page else (sample_count + rows * cols - 1) // (rows * cols),
        'single_page': single_page
    }
